# -*- coding: utf-8 -*-
# ===============================================================
#
#    Copyright (C) 2024 Beike, Inc. All Rights Reserved.
#
#    @Create Author : luxu
#    @Create Time   : 2024/7/16
#    @Description   : 
#
# ===============================================================
import io
import logging

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE

from doc_parser.layout_parser.layout.simple_block import SimpleBlock
from services.constants import TEXT, TABLE, IMAGE
from services.layout_parse_utils import get_s3_links_for_simple_block_batch, trans_simple_block_list2string


def layout_parse(file):
    simple_block_list = []

    pptx_stream = io.BytesIO(file)
    pr = Presentation(pptx_stream)
    # 遍历幻灯片中的所有形状
    for slide in pr.slides:
        for shape in slide.shapes:
            try:
                # placeholder需要识别
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:  # 14
                    placeholder_type = shape.placeholder_format.type
                    if placeholder_type == PP_PLACEHOLDER_TYPE.PICTURE:  # 18
                        simple_block_list.append(SimpleBlock(type=IMAGE, image_bytes=shape.image.blob))
                    else:
                        # 文字占位符
                        simple_block_list.append(SimpleBlock(type=TEXT, text=""))
                # 图片
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 13
                    try:
                        simple_block_list.append(SimpleBlock(type=IMAGE, image_bytes=shape.image.blob))
                    except Exception as e:
                        simple_block_list.append(SimpleBlock(type=IMAGE, image_bytes=None))
                # 文字
                # todo ppt中组合可能无法解析
                elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:  # 17
                    if shape.text:
                        simple_block_list.append(SimpleBlock(type=TEXT, text=shape.text))
                # 表格
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # 19
                    table = shape.table
                    table_text = ""
                    for row in table.rows:
                        for cell in row.cells:
                            table_text = " | ".join([table_text, cell.text])
                    if table_text:
                        simple_block_list.append(SimpleBlock(type=TABLE, text=table_text))

            except Exception as e:
                logging.error(f"处理元素时出错，type: {shape.shape_type}，errmsg：{str(e)}")
                continue

    # SimpleBlock的list批量获取S3链接，并返回目标结构
    result_json = get_s3_links_for_simple_block_batch(simple_block_list)
    result_text = trans_simple_block_list2string(result_json)
    return result_json, result_text


def get_page_count(file):
    pptx_stream = io.BytesIO(file)
    pr = Presentation(pptx_stream)
    slide_count = len(pr.slides)
    return slide_count