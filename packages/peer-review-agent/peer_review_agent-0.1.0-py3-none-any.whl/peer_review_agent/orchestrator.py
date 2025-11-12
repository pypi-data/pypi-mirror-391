"""
Core orchestration logic for the Peer Review Agent CLI.

The module exposes two public helpers:
    - `process_articles(...)` – batch processing entry point.
    - `process_article(...)` – single-file processing (also used internally).
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.util import Pt
from PyPDF2 import PdfReader

try:
    from docx import Document
    from docx.shared import RGBColor
except ImportError as exc:  # pragma: no cover
    raise ImportError(
        "python-docx is required for peer-review outputs. Install via `pip install python-docx`."
    ) from exc


SUPPORTED_EXTS = (".pdf", ".docx")
DERIVED_MARKERS = ("_auto_", "_peer_review", "_annotated", "_redline")


@dataclass
class ArticleMetrics:
    """Container for heuristic signals extracted from the manuscript text."""

    title: str
    authors_line: str
    design: str
    sample_size: Optional[int]
    follow_up_notes: List[str]
    loss_to_follow_up: bool
    stat_methods: List[str]
    multivariable: bool
    power_reported: bool
    interactions_reported: bool
    limitations_reported: bool
    discussion_present: bool
    precision_terms: bool
    external_validity_flags: List[str] = field(default_factory=list)


DESIGN_KEYWORDS: List[Tuple[str, str]] = [
    ("randomized", "Randomized controlled trial"),
    ("randomised", "Randomized controlled trial"),
    ("prospective", "Prospective cohort"),
    ("retrospective", "Retrospective cohort"),
    ("case-control", "Case-control study"),
    ("cross-sectional", "Cross-sectional study"),
    ("meta-analysis", "Meta-analysis"),
]

STAT_METHODS: List[Tuple[str, str]] = [
    ("poisson", "Poisson regression"),
    ("cox", "Cox proportional hazards"),
    ("logistic", "Logistic regression"),
    ("hazard ratio", "Hazard ratios reported"),
    ("odds ratio", "Odds ratios reported"),
    ("kaplan-meier", "Kaplan-Meier analysis"),
    ("mixed model", "Mixed-effects models"),
    ("anova", "ANOVA / general linear models"),
    ("chi-square", "Chi-square tests"),
    ("multivariate", "Multivariable modeling"),
    ("generalized linear", "Generalized linear models"),
]


# --------------------------------------------------------------------------- #
# Text extraction helpers


def _read_pdf_text(pdf_path: Path) -> str:
    reader = PdfReader(str(pdf_path))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _read_docx_text(docx_path: Path) -> str:
    doc = Document(str(docx_path))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def _extract_text(article_path: Path) -> str:
    suffix = article_path.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf_text(article_path)
    if suffix == ".docx":
        return _read_docx_text(article_path)
    raise ValueError(f"Unsupported file type: {article_path}")


# --------------------------------------------------------------------------- #
# Heuristic detection helpers


def _first_non_empty(lines: Sequence[str]) -> str:
    for line in lines:
        stripped = line.strip()
        if stripped:
            return stripped
    return "Untitled Article"


def _second_non_empty(lines: Sequence[str]) -> str:
    non_empty = [ln.strip() for ln in lines if ln.strip()]
    return non_empty[1] if len(non_empty) > 1 else ""


def _detect_design(text_lower: str) -> str:
    for keyword, label in DESIGN_KEYWORDS:
        if keyword in text_lower:
            return label
    return "Design not explicitly stated"


def _detect_sample_size(text: str) -> Optional[int]:
    hits: List[int] = []
    for match in re.finditer(
        r"(?:n\s*=\s*)?([0-9][0-9,]{1,7})\s+(participants|patients|subjects|cases|individuals)",
        text,
        flags=re.IGNORECASE,
    ):
        hits.append(int(match.group(1).replace(",", "")))
    for match in re.finditer(r"\bN\s*=\s*([0-9][0-9,]{1,7})", text):
        hits.append(int(match.group(1).replace(",", "")))
    return max(hits) if hits else None


def _sentences_containing(text: str, keyword: str) -> List[str]:
    sentences = re.split(r"(?<=[.!?])\s+", text)
    return [s.strip() for s in sentences if keyword.lower() in s.lower()]


def _detect_stat_methods(text_lower: str) -> List[str]:
    return sorted({label for keyword, label in STAT_METHODS if keyword in text_lower})


def _build_metrics(text: str) -> ArticleMetrics:
    lines = text.splitlines()
    text_lower = text.lower()
    title = _first_non_empty(lines)
    authors_line = _second_non_empty(lines)
    design = _detect_design(text_lower)
    sample_size = _detect_sample_size(text)
    follow_up_notes = _sentences_containing(text, "follow-up")
    loss_to_follow_up = "loss to follow-up" in text_lower or "lost to follow-up" in text_lower
    stat_methods = _detect_stat_methods(text_lower)
    multivariable = any(token in text_lower for token in ("multivariable", "multivariate", "adjusted"))
    power_reported = "power" in text_lower and "sample size" in text_lower
    interactions_reported = "interaction" in text_lower
    limitations_reported = "limitation" in text_lower
    discussion_present = "discussion" in text_lower
    precision_terms = "confidence interval" in text_lower or "95%" in text_lower or "ci " in text_lower
    external_flags: List[str] = []
    if "tertiary" in text_lower:
        external_flags.append("Tertiary-care recruitment noted.")
    if "single center" in text_lower:
        external_flags.append("Single-center design referenced.")
    if "multicenter" in text_lower or "multi-center" in text_lower:
        external_flags.append("Multicenter sampling referenced.")
    return ArticleMetrics(
        title=title,
        authors_line=authors_line,
        design=design,
        sample_size=sample_size,
        follow_up_notes=follow_up_notes,
        loss_to_follow_up=loss_to_follow_up,
        stat_methods=stat_methods,
        multivariable=multivariable,
        power_reported=power_reported,
        interactions_reported=interactions_reported,
        limitations_reported=limitations_reported,
        discussion_present=discussion_present,
        precision_terms=precision_terms,
        external_validity_flags=external_flags,
    )


# --------------------------------------------------------------------------- #
# Narrative builders


def _strengths_from_metrics(metrics: ArticleMetrics) -> List[str]:
    strengths: List[str] = []
    if metrics.sample_size:
        descriptor = "Large sample" if metrics.sample_size >= 1000 else "Sample size captured"
        strengths.append(f"{descriptor} (≈{metrics.sample_size}).")
    if metrics.follow_up_notes:
        strengths.append("Follow-up procedures described, enabling longitudinal outcomes.")
    if metrics.multivariable:
        strengths.append("Adjusted/ multivariable modeling acknowledged.")
    if metrics.stat_methods:
        strengths.append("Statistical approaches reported: " + ", ".join(metrics.stat_methods) + ".")
    if metrics.precision_terms:
        strengths.append("Confidence intervals or precision metrics provided.")
    if metrics.limitations_reported:
        strengths.append("Manuscript includes a limitations discussion.")
    return strengths


def _gaps_from_metrics(metrics: ArticleMetrics) -> List[str]:
    gaps: List[str] = []
    if not metrics.power_reported:
        gaps.append("No sample size/ power justification detected.")
    if not metrics.interactions_reported:
        gaps.append("Interaction tests/effect modifiers not described.")
    if metrics.design != "Randomized controlled trial":
        gaps.append("Non-randomized design susceptible to confounding.")
    if not metrics.follow_up_notes:
        gaps.append("Follow-up duration/process not documented.")
    if metrics.loss_to_follow_up:
        gaps.append("Loss to follow-up noted; attrition handling unclear.")
    if not metrics.precision_terms:
        gaps.append("Precision metrics (CI/SE) absent.")
    if not metrics.limitations_reported:
        gaps.append("Limitations section missing or insufficient.")
    return gaps


def _render_report(
    metrics: ArticleMetrics,
    source: Path,
    text_path: Path,
    strengths: List[str],
    gaps: List[str],
) -> str:
    stat_methods = metrics.stat_methods or ["Not specified"]
    follow_up = metrics.follow_up_notes[:2] if metrics.follow_up_notes else ["Not described."]
    external = metrics.external_validity_flags or ["No explicit recruitment context detected."]
    lines: List[str] = []
    lines.append(f"# Automated Critical Analysis: {metrics.title}")
    lines.append("")
    lines.append("## Article Snapshot")
    lines.append(f"- **Source file:** {source.name}")
    lines.append(f"- **Title line:** {metrics.title}")
    lines.append(f"- **Authors line:** {metrics.authors_line or 'Not detected'}")
    lines.append(f"- **Design cue:** {metrics.design}")
    lines.append(f"- **Sample size (heuristic):** {metrics.sample_size or 'Not detected'}")
    lines.append("")
    lines.append("## Automated Strength Signals")
    if strengths:
        lines.extend(f"- {signal}" for signal in strengths)
    else:
        lines.append("- No automatic strengths detected; review manually.")
    lines.append("")
    lines.append("## Critical Gaps & Risks")
    if gaps:
        lines.extend(f"- {gap}" for gap in gaps)
    else:
        lines.append("- No automatic gaps detected; confirm manually.")
    lines.append("")
    lines.append("## Sample Size, Power, and Precision")
    lines.append(f"- Sample size flag: {metrics.sample_size or 'Not parsed from text'}.")
    lines.append(f"- Power calculation mentioned: {'Yes' if metrics.power_reported else 'No'}")
    lines.append(f"- Precision metrics detected: {'Yes' if metrics.precision_terms else 'No'}")
    lines.append("")
    lines.append("## Analytical Methods Assessment")
    lines.append(f"- Statistical techniques referenced: {', '.join(stat_methods)}.")
    lines.append(f"- Multivariable modeling noted: {'Yes' if metrics.multivariable else 'Not detected'}")
    lines.append(f"- Interaction testing noted: {'Yes' if metrics.interactions_reported else 'Not mentioned'}")
    lines.append("")
    lines.append("## Follow-up and Outcome Assessment")
    for sentence in follow_up:
        lines.append(f"- {sentence}")
    lines.append(f"- Loss to follow-up references detected: {'Yes' if metrics.loss_to_follow_up else 'No'}")
    lines.append("")
    lines.append("## Discussion, Bias, and Validity Considerations")
    lines.append(f"- Discussion section detected: {'Yes' if metrics.discussion_present else 'No'}")
    lines.append(f"- Limitations discussed: {'Yes' if metrics.limitations_reported else 'No'}")
    lines.append(f"- External validity cues: {'; '.join(external)}")
    lines.append("")
    lines.append("## Automation Notes")
    lines.append(f"- Raw text stored at: {text_path.name}")
    lines.append("- Heuristic outputs depend on textual cues; verify critical details manually.")
    return "\n".join(lines)


# --------------------------------------------------------------------------- #
# Output generators


def _generate_ppt(metrics: ArticleMetrics, strengths: List[str], gaps: List[str], ppt_path: Path) -> None:
    prs = Presentation()

    def add_bullets(title: str, bullets: Iterable[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(20)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = metrics.title[:120]
    subtitle_bits = [metrics.design, f"Sample size: {metrics.sample_size or 'n/a'}"]
    title_slide.placeholders[1].text = " | ".join(bit for bit in subtitle_bits if bit)

    add_bullets(
        "Study Snapshot",
        [
            f"Design cue: {metrics.design}",
            f"Sample size (detected): {metrics.sample_size or 'Not found'}",
            f"Statistical methods: {', '.join(metrics.stat_methods) if metrics.stat_methods else 'Not specified'}",
        ],
    )
    add_bullets("Automated Strengths", strengths or ["No strength cues detected automatically."])
    add_bullets("Critical Gaps", gaps or ["No gaps detected automatically; manual audit required."])
    add_bullets(
        "Sample Size & Power",
        [
            f"Sample size: {metrics.sample_size or 'n/a'}",
            f"Power mentioned: {'Yes' if metrics.power_reported else 'No'}",
            f"Precision metrics present: {'Yes' if metrics.precision_terms else 'No'}",
        ],
    )
    add_bullets(
        "Analytical Methods",
        [
            f"Multivariable modeling: {'Yes' if metrics.multivariable else 'Not detected'}",
            f"Interaction testing: {'Yes' if metrics.interactions_reported else 'Not mentioned'}",
            f"Loss to follow-up cited: {'Yes' if metrics.loss_to_follow_up else 'No'}",
        ],
    )
    add_bullets(
        "Bias & Validity",
        [
            f"Discussion present: {'Yes' if metrics.discussion_present else 'No'}",
            f"Limitations noted: {'Yes' if metrics.limitations_reported else 'No'}",
            f"External validity cues: {', '.join(metrics.external_validity_flags) or 'Not detected'}",
        ],
    )
    add_bullets(
        "Recommended Next Checks",
        [
            "Verify strengths/gaps against the full manuscript.",
            "Confirm sample size & power statements.",
            "Add local clinical context or visuals before presenting.",
        ],
    )

    prs.save(ppt_path)


def _build_track_change_suggestions(metrics: ArticleMetrics, text: str) -> List[dict]:
    suggestions: List[dict] = []
    if not metrics.power_reported:
        excerpt = _sentences_containing(text, "sample size")
        suggestions.append(
            {
                "comment": "Provide a formal sample size/ power justification in Methods.",
                "excerpt": excerpt[0] if excerpt else "",
            }
        )
    if metrics.loss_to_follow_up:
        excerpt = _sentences_containing(text, "loss to follow-up")
        suggestions.append(
            {
                "comment": "Clarify handling of loss to follow-up (sensitivity analysis or imputation).",
                "excerpt": excerpt[0] if excerpt else "",
            }
        )
    if not metrics.interactions_reported:
        suggestions.append(
            {
                "comment": "Assess potential effect modifiers (e.g., sex × BMI) and report interaction tests.",
                "excerpt": "",
            }
        )
    if metrics.follow_up_notes:
        suggestions.append(
            {
                "comment": "Condense follow-up logistics into a figure/table for clarity.",
                "excerpt": metrics.follow_up_notes[0],
            }
        )
    return suggestions


def _generate_peer_review_doc(
    metrics: ArticleMetrics,
    strengths: List[str],
    gaps: List[str],
    text: str,
    output_path: Path,
) -> None:
    doc = Document()
    doc.add_heading("Peer Review & Critical Comments", level=1)
    doc.add_paragraph(f"Article: {metrics.title}")
    doc.add_paragraph(f"Design cue: {metrics.design}")
    doc.add_paragraph(f"Detected sample size: {metrics.sample_size or 'Not parsed'}")

    doc.add_heading("Major Comments", level=2)
    if gaps:
        for idx, comment in enumerate(gaps, start=1):
            doc.add_paragraph(f"{idx}. {comment}")
    else:
        doc.add_paragraph("No major risks detected automatically.")

    doc.add_heading("Minor Comments / Positives", level=2)
    if strengths:
        for idx, signal in enumerate(strengths, start=1):
            doc.add_paragraph(f"{idx}. {signal}")
    else:
        doc.add_paragraph("No automated strengths captured.")

    suggestions = _build_track_change_suggestions(metrics, text)
    doc.add_heading("Suggested Track Changes / Inline Comments", level=2)
    if suggestions:
        for idx, suggestion in enumerate(suggestions, start=1):
            doc.add_paragraph(f"{idx}. {suggestion['comment']}")
            if suggestion["excerpt"]:
                excerpt_para = doc.add_paragraph(f"Excerpt: {suggestion['excerpt']}")
                excerpt_para.italic = True
    else:
        doc.add_paragraph("No suggestions generated automatically.")

    doc.save(output_path)


def _generate_annotation_text(text: str, suggestions: List[dict], output_path: Path) -> None:
    annotated = text
    for idx, suggestion in enumerate(suggestions, start=1):
        excerpt = suggestion.get("excerpt")
        note = f" [[Comment {idx}: {suggestion['comment']}]]"
        if excerpt and excerpt in annotated:
            annotated = annotated.replace(excerpt, f"{excerpt}{note}", 1)
        else:
            annotated += f"\n\n[Comment {idx}] {suggestion['comment']}"
    output_path.write_text(annotated, encoding="utf-8")


def _derive_action_items(metrics: ArticleMetrics) -> List[str]:
    items = [
        "Ensure the abstract states design, sample size, and quantitative results with 95% confidence intervals.",
        "Explicitly document statistical tests, regression models, and software used.",
        "Align conclusions with observational design; avoid causal claims without randomized data.",
    ]
    if not metrics.power_reported:
        items.append("Add a sample size/ power justification in Methods.")
    if not metrics.precision_terms:
        items.append("Report effect sizes with 95% confidence intervals.")
    if not metrics.interactions_reported:
        items.append("Document whether interaction/effect-modifier analyses were performed.")
    if metrics.loss_to_follow_up:
        items.append("Describe handling of loss to follow-up or conduct sensitivity analyses.")
    if not metrics.limitations_reported:
        items.append("Include a dedicated Limitations section discussing selection bias and missing data.")
    return items


def _derive_section_guidance(metrics: ArticleMetrics) -> Dict[str, List[Tuple[str, str]]]:
    guidance: Dict[str, List[Tuple[str, str]]] = {
        "ABSTRACT": [
            (
                "Clarify study design, sample size, and primary outcomes with quantitative values.",
                "Suggested rewrite: 'Methods: Retrospective cohort of 2006 drug-susceptible TB patients; Poisson models estimated adjusted incidence rate ratios with 95% CI.'",
            )
        ],
        "INTRODUCTION": [
            (
                "Focus on the specific knowledge gap this manuscript addresses rather than broad epidemiology.",
                "Suggested rewrite: 'Despite national TB surveillance, few multicenter cohorts report post-treatment recurrence under programmatic conditions; this study addresses that gap.'",
            )
        ],
        "METHODS": [
            (
                "Detail sampling frame, inclusion/exclusion criteria, assay platforms, and data sources.",
                "Suggested rewrite: 'We consecutively enrolled adults ≥15 years with microbiologically confirmed TB, excluding MDR cases and records lacking HIV results; thyroid assays used chemiluminescent immunoassay (CV<5%).'",
            )
        ],
        "STATISTICAL ANALYSIS": [
            (
                "Specify statistical tests, regression models, covariate selection, and software versions.",
                "Suggested rewrite: 'Comparisons used chi-square or Fisher exact tests; multivariable Poisson regression with site-level clustering generated adjusted IRRs (Stata 17).'",
            )
        ],
        "RESULTS": [
            (
                "Provide numerators/denominators and 95% CIs for key outcomes rather than narrative statements.",
                "Suggested rewrite: 'Unfavorable outcomes occurred in 365/2006 (18.2%, 95% CI 16.5-19.9); loss to follow-up accounted for 137 cases (6.8%).'",
            )
        ],
        "DISCUSSION": [
            (
                "Interpret findings relative to similar cohorts and discuss plausible mechanisms for discrepancies.",
                "Suggested rewrite: 'Higher loss to follow-up versus RePORT Brazil likely reflects inpatient recruitment and limited post-discharge tracing at our sites.'",
            )
        ],
        "CONCLUSION": [
            (
                "Add a Limitations paragraph acknowledging retrospective design, tertiary sampling, and missing adherence data.",
                "Suggested rewrite: 'Limitations include retrospective abstraction, reliance on tertiary centers, and absence of adherence/ glycemic control data, which may limit generalizability.'",
            )
        ],
        "REFERENCES": [
            (
                "Verify citation formatting and ensure all guidelines cited in text appear in the reference list.",
                "Suggested rewrite: 'Add WHO 2023 End TB report and NTEP 2025 guidelines referenced in the introduction.'",
            )
        ],
    }

    if not metrics.power_reported:
        guidance["METHODS"].append(
            (
                "Include a sample size/ power justification referencing prior prevalence or effect estimates.",
                "Suggested rewrite: 'Assuming 20% prevalence with 5% precision (alpha=0.05), minimum n=246; we enrolled 2006 participants to enable site-level analyses.'",
            )
        )
    if not metrics.interactions_reported:
        guidance["STATISTICAL ANALYSIS"].append(
            (
                "State whether interaction/ effect-modifier analyses (e.g., sex × BMI) were performed or justify omission.",
                "Suggested rewrite: 'We assessed sex × BMI interaction; non-significant terms (p>0.10) were removed from final models.'",
            )
        )
    if metrics.loss_to_follow_up:
        guidance["RESULTS"].append(
            (
                "Discuss predictors of loss to follow-up and whether sensitivity analyses were performed.",
                "Suggested rewrite: 'Logistic regression identified underweight males as high risk for loss to follow-up (adjusted OR 1.9, 95% CI 1.3-2.6).'",
            )
        )
    return guidance


def _normalize_section_key(text: str) -> Optional[str]:
    cleaned = text.strip().lower()
    if not cleaned:
        return None
    aliases = {
        "abstract": "ABSTRACT",
        "introduction": "INTRODUCTION",
        "materials and methods": "METHODS",
        "methods": "METHODS",
        "statistical analysis": "STATISTICAL ANALYSIS",
        "analysis": "STATISTICAL ANALYSIS",
        "results": "RESULTS",
        "discussion": "DISCUSSION",
        "conclusion": "CONCLUSION",
        "conclusions": "CONCLUSION",
        "references": "REFERENCES",
    }
    return aliases.get(cleaned)


def _generate_redline_doc(
    article_path: Path,
    raw_text: str,
    metrics: ArticleMetrics,
    action_items: List[str],
    section_guidance: Dict[str, List[Tuple[str, str]]],
    output_path: Path,
) -> None:
    redline_doc = Document()
    redline_doc.add_heading("Peer Review Redline + Action Plan", level=1)
    redline_doc.add_paragraph(
        "Reviewer summary: Manuscript requires revisions to methodological transparency, "
        "statistical reporting, and conclusion alignment."
    )
    redline_doc.add_heading("Prioritized Action List", level=2)
    for idx, item in enumerate(action_items, start=1):
        redline_doc.add_paragraph(f"{idx}. {item}")
    redline_doc.add_heading("Annotated Manuscript", level=2)

    def copy_paragraph(text: str) -> None:
        redline_doc.add_paragraph(text)

    def add_comment(text: str) -> None:
        para = redline_doc.add_paragraph()
        run = para.add_run(text)
        run.font.color.rgb = RGBColor(0xB0, 0x00, 0x00)
        run.font.bold = True

    def add_suggestion(text: str) -> None:
        para = redline_doc.add_paragraph()
        run = para.add_run(text)
        run.font.color.rgb = RGBColor(0x00, 0x45, 0x8A)
        run.font.italic = True

    remaining_guidance = {key: list(value) for key, value in section_guidance.items()}

    if article_path.suffix.lower() == ".docx":
        src_doc = Document(str(article_path))
        paragraphs = [para.text for para in src_doc.paragraphs]
    else:
        paragraphs = [line.strip() for line in raw_text.splitlines() if line.strip()]

    for paragraph in paragraphs:
        copy_paragraph(paragraph)
        section_key = _normalize_section_key(paragraph)
        if section_key and remaining_guidance.get(section_key):
            for comment, suggestion in remaining_guidance[section_key]:
                add_comment(comment)
                add_suggestion(suggestion)
            remaining_guidance[section_key] = []

    redline_doc.save(output_path)


# --------------------------------------------------------------------------- #
# Public processing API


def process_article(
    article_path: Path,
    *,
    force: bool = False,
    generate_peer_review: bool = False,
    annotate: bool = False,
    redline: bool = False,
) -> None:
    """Process a single article file."""
    lowered = article_path.name.lower()
    if any(marker in lowered for marker in DERIVED_MARKERS):
        print(f"[skip] Derived artifact skipped: {article_path.name}")
        return

    folder = article_path.parent
    stem = article_path.stem.replace(" ", "_")
    text_path = folder / f"{stem}.txt"
    report_path = folder / f"{stem}_auto_critical_analysis.md"
    ppt_path = folder / f"{stem}_auto_appraisal.pptx"
    peer_path = folder / f"{stem}_peer_review.docx"
    annotations_path = folder / f"{stem}_annotated_comments.txt"
    redline_path = folder / f"{stem}_redline_review.docx"

    if (
        not force
        and report_path.exists()
        and ppt_path.exists()
        and (not generate_peer_review or peer_path.exists())
        and (not annotate or annotations_path.exists())
        and (not redline or redline_path.exists())
    ):
        print(f"[skip] Outputs already exist for {article_path.name}")
        return

    print(f"[info] Processing {article_path.name}")
    text = _extract_text(article_path)
    text_path.write_text(text, encoding="utf-8")
    metrics = _build_metrics(text)
    strengths = _strengths_from_metrics(metrics)
    gaps = _gaps_from_metrics(metrics)
    report = _render_report(metrics, article_path, text_path, strengths, gaps)
    report_path.write_text(report, encoding="utf-8")
    _generate_ppt(metrics, strengths, gaps, ppt_path)

    suggestions = _build_track_change_suggestions(metrics, text)
    if generate_peer_review:
        _generate_peer_review_doc(metrics, strengths, gaps, text, peer_path)
    if annotate:
        _generate_annotation_text(text, suggestions, annotations_path)
    if redline:
        action_items = _derive_action_items(metrics)
        section_guidance = _derive_section_guidance(metrics)
        _generate_redline_doc(article_path, text, metrics, action_items, section_guidance, redline_path)

    print(f"[done] Generated outputs for {article_path.name}")


def _discover_targets(root: Path) -> List[Path]:
    folders: List[Path] = []
    for child in root.iterdir():
        if child.is_dir():
            folders.append(child)
    return sorted(folders)


def process_articles(
    *,
    root: Path,
    folder: Optional[Path],
    force: bool,
    peer_review: bool,
    annotate: bool,
    redline: bool,
) -> None:
    """Batch processing entry point."""
    targets = [folder] if folder else _discover_targets(root)
    if not targets:
        print("[warn] No target folders discovered.")
        return

    for target in targets:
        if not target.exists():
            print(f"[warn] Folder not found: {target}")
            continue
        article_files = [
            path for path in target.iterdir() if path.suffix.lower() in SUPPORTED_EXTS
        ]
        if not article_files:
            continue
        for article in article_files:
            process_article(
                article,
                force=force,
                generate_peer_review=peer_review,
                annotate=annotate,
                redline=redline,
            )

