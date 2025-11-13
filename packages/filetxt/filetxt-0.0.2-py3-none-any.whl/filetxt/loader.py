import os
import codecs
from datetime import datetime
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_community.document_loaders.text import TextLoader
from docx import Document
from ebooklib import epub
from bs4 import BeautifulSoup
from pptx import Presentation
from openpyxl import load_workbook
import mobi
from pypdf import PdfReader
import html2text

class BaseLoader:
	def __init__(self, file_path):
		self.filepath = file_path

	def load(self):
		pts = []
		for i, t in self.load_pagetext():
			pts.append(' '.join(t.split('\t')))
		return '\n'.join(pts)
	
	def load_pagetext(self):
		raise Exception('Not implement')
			
class MyMobiLoader(BaseLoader):
	def load_pagetext(self):
		tempdir, filepath = mobi.extract(self.filepath)
		with codecs.open(filepath, "r", "utf-8") as f:
			content=f.read()
			yield (0, html2text.html2text(content))
		
class MyPdfLoader(BaseLoader):
	def load_pagetext(self):
		"""Reads the PDF file and returns all text as a single string."""
		reader = PdfReader(self.filepath)
		pts = []
		for i,page in enumerate(reader.pages):
			t = page.extract_text() or ''
			yield (i, t)
	
class MyDocxLoader(BaseLoader):
	def load_pagetext(self):
		"""Reads the .docx file and returns the full text as a single string."""
		docx = Document(self.filepath)
		for i, para in enumerate(docx.paragraphs):
			yield (i,para.text) 

class MyPptLoader(BaseLoader):
	def load_pagetext(self):
		prs = Presentation(self.filepath)
		for i, slide in enumerate(prs.slides):
			txts = []
			for shape in slide.shapes:
				if hasattr(shape, "text"):
					txts.append(shape.text)
			yield (i,'\n'.join(txts))

class MyCsvLoader(BaseLoader):
	def load_pagetext(self):
		loader = CSVLoader(self.filepath)
		docs = loader.load()
		for i, d in enumerate(docs):
			dat = (i, d.page_content) 
			yield dat

class MyExcelLoader(BaseLoader):
	def load_pagetext(self):
		"""Reads all sheets in the Excel file and returns the content as a string."""
		self.workbook = load_workbook(filename=self.filepath, data_only=True)
		content = []

		for i, sheet in enumerate(self.workbook.worksheets):
			txts = []
			for row in sheet.iter_rows(values_only=True):
				row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
				txts.append(row_text)
			yield(i, '\n'.join(txts))

		return content

class MyEpubLoader(BaseLoader):
	def __init__(self, file_path):
		self.filepath = file_path
		self.book = None

	def load_pagetext(self):
		"""Reads the EPUB file and returns all text content as a string."""
		self.book = epub.read_epub(self.filepath)
		for i, item in enumerate(self.book.get_items()):
			if isinstance(item, epub.EpubHtml):
				soup = BeautifulSoup(item.get_content(), 'html.parser')
				text = soup.get_text()
				yield(i, text.strip())


class MyTextLoader(BaseLoader):
	def load_pagetext(self):
		loader = TextLoader(self.filepath)
		docs = loader.load()
		for i, d in enumerate(docs):
			dat = (i, d.page_content) 
			yield dat

class File2Text:
	all_loaders = {
		'docx':MyDocxLoader,
		'pptx':MyPptLoader,
		'csv':MyCsvLoader,
		'xlsx':MyExcelLoader,
		'pdf':MyPdfLoader,
		'epub':MyEpubLoader,
		'mobi':MyMobiLoader,
		'md': MyTextLoader,
		'txt':MyTextLoader
	}
	@classmethod
	def supported_types(self):
		return [k for k in self.all_loaders.keys()]

	def __init__(self, filepath):
		self.filepath = filepath

	def load_pagetext(self):
		k = self.filepath.lower().split('.')[-1]
		klass = self.all_loaders.get(k, MyTextLoader)
		loader = klass(self.filepath)
		for d in loader.load_pagetext():
			yield d

	def load(self):
		k = self.filepath.lower().split('.')[-1]
		klass = self.all_loaders.get(k, MyTextLoader)
		loader = klass(self.filepath)
		return loader.load()
			
def fileloader(file_path):
	loader = File2Text(file_path)
	return loader.load()

def filepageloader(file_path):
	loader = File2Text(file_path)
	for d in loader.load_pagetext():
		yield d

if __name__ == '__main__':
	import sys
	if len(sys.argv) < 2:
		print(f'{sys.argv[0]} file\nload a file and get its text')
		sys.exit(1)
	text = fileloader(sys.argv[1])
	print(f'{text=}')
	"""
	for i, txt in filepageloader(sys.argv[1]):
		print(f'page:{i}\n{txt}\n=======\n')
	"""
