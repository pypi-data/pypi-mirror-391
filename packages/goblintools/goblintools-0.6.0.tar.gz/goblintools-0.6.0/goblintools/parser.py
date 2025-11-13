import csv
import logging
import os
from typing import Dict, Callable, Optional, Union, List
from pathlib import Path
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from dbfread import DBF
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from pypdf.generic import IndirectObject
import openpyxl
import xlrd
from odf import text, teletype
from odf.opendocument import load
from odf.text import P
import docx
from functools import lru_cache
from goblintools.ocr_parser import OCRProcessor
from goblintools.config import GoblinConfig, OCRConfig

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TextExtractor:
    """Main class for handling text extraction from various file formats."""

    def __init__(self, ocr_handler=False, use_aws=False, aws_access_key=None, aws_secret_key=None, aws_region='us-east-1', config: Optional[GoblinConfig] = None):
        """
        Initialize the text extractor.

        Args:
            ocr_handler: Enable OCR for image-based PDFs
            use_aws: Use AWS Textract for OCR
            aws_access_key: AWS access key
            aws_secret_key: AWS secret key
            aws_region: AWS region
            config: GoblinConfig object (overrides other parameters)
        """
        self.config = config or GoblinConfig.default()

        # Override config with explicit parameters if provided
        if any([use_aws, aws_access_key, aws_secret_key, aws_region != 'us-east-1']):
            self.config.ocr = OCRConfig(use_aws, aws_access_key, aws_secret_key, aws_region)

        if ocr_handler:
            self.ocr_handler = OCRProcessor(self.config.ocr)
        else:
            self.ocr_handler = None

        self._parsers = None  # Lazy initialization

    @property
    def parsers(self) -> Dict[str, Callable]:
        """Lazy-loaded parsers dictionary"""
        if self._parsers is None:
            self._parsers = self._initialize_parsers()
        return self._parsers

    def _initialize_parsers(self) -> Dict[str, Callable]:
        """Initialize all available text extraction parsers."""
        return {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.txt': self._extract_txt,
            '.pptx': self._extract_pptx,
            '.html': self._extract_html,
            '.odt': self._extract_odt,
            '.rtf': self._extract_rtf,
            '.csv': self._extract_csv,
            '.xml': self._extract_xml,
            '.xlsx': self._extract_xlsx,
            '.xlsm': self._extract_xlsx,
            '.xls': self._extract_xls,
            '.ods': self._extract_ods,
            '.dbf': self._extract_dbf,
        }

    def add_parser(self, extension: str, parser_func: Callable) -> None:
        """Add or override a parser for a specific file extension."""
        self.parsers[extension.lower()] = parser_func

    def _extract_with_metadata(self, file_path: str, parser: Callable) -> Dict:
        """Extract text with metadata structure."""
        filename = Path(file_path).name
        file_extension = Path(file_path).suffix.lower()

        if file_extension == '.pdf':
            return self._extract_pdf_with_metadata(file_path)
        else:
            text = parser(file_path)
            if text:
                metadata = {filename: {1: text}}
                return {"text": text, "metadata": metadata}
            else:
                return {"text": "", "metadata": {}}

    def extract_from_file(self, file_path: str, include_metadata: bool = False) -> Union[str, Dict]:
        """
        Extract text from a single file.

        Args:
            file_path: Path to the file to extract text from
            include_metadata: If True, return dict with text and metadata

        Returns:
            Extracted text as string or dict with text and metadata
        """
        if not os.path.exists(file_path):
            logger.warning(f"File not found: {file_path}")
            if include_metadata:
                return {"text": "", "metadata": {}}
            return ""

        file_extension = Path(file_path).suffix.lower()
        parser = self.parsers.get(file_extension)

        if not parser:
            logger.warning(f"No parser available for file extension: {file_extension}")
            if include_metadata:
                return {"text": "", "metadata": {}}
            return ""

        try:
            if include_metadata:
                return self._extract_with_metadata(file_path, parser)
            else:
                return parser(file_path)
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            if include_metadata:
                return {"text": "", "metadata": {}}
            return ""

    def extract_from_folder(self, folder_path: str, include_metadata: bool = False) -> Union[str, Dict]:
        """
        Extract text from all supported files in a folder (recursively).

        Args:
            folder_path: Path to the folder to process
            include_metadata: If True, return dict with combined text and metadata

        Returns:
            Combined extracted text or dict with text and metadata
        """
        if not os.path.exists(folder_path):
            logger.warning(f"Folder not found: {folder_path}")
            if include_metadata:
                return {"text": "", "metadata": {}}
            return ""

        if include_metadata:
            all_texts = []
            all_metadata = {}

            for root, _, files in os.walk(folder_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    try:
                        if os.path.getsize(file_path) > self.config.max_file_size:
                            logger.warning(f"Skipping large file: {file_path}")
                            continue
                    except OSError:
                        continue

                    result = self.extract_from_file(file_path, include_metadata=True)
                    if result and result.get("text") and result.get("metadata"):
                        all_texts.append(result["text"])
                        all_metadata.update(result["metadata"])

            if not all_texts:
                return {"text": "", "metadata": {}}

            return {
                "text": " ".join(all_texts),
                "metadata": all_metadata
            }
        else:
            def text_generator():
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        file_path = os.path.join(root, file)
                        try:
                            if os.path.getsize(file_path) > self.config.max_file_size:
                                logger.warning(f"Skipping large file: {file_path}")
                                continue
                        except OSError:
                            continue

                        text = self.extract_from_file(file_path)
                        if text:
                            yield text
            return ' '.join(text_generator())

    def pdf_needs_ocr(self, file_path: str) -> bool:
        """Check if PDF needs OCR processing"""
        try:
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text and not text.isspace():
                        return False
            return True
        except Exception as e:
            logger.error(f"Error checking PDF {file_path}: {e}")
            return True

    def _resave_pdf(self, file_path: str) -> str:
        """Resave PDF to fix potential issues"""
        reader = PdfReader(file_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        output_path = Path(file_path).with_suffix(".resaved.pdf")
        with open(output_path, 'wb') as f:
            writer.write(f)

        return str(output_path)

    def validate_installation(self) -> Dict[str, bool]:
        """Check if all dependencies are properly installed"""
        results = {}

        # Check Tesseract
        try:
            import pytesseract
            pytesseract.get_tesseract_version()
            results['tesseract'] = True
        except:
            results['tesseract'] = False

        # Check AWS credentials
        if self.ocr_handler and self.config.ocr.use_aws:
            try:
                self.ocr_handler.textract_client.list_document_analysis_jobs
                results['aws_textract'] = True
            except:
                results['aws_textract'] = False

        return results

    # Individual parser methods
    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF files using PyPDF, with fallback to OCR if needed."""
        extracted_text = ''
        has_images = False
        temp_file = None

        try:
            temp_file = self._resave_pdf(file_path)
            reader = PdfReader(temp_file)

            for i, page in enumerate(reader.pages):
                try:
                    extracted_text += page.extract_text() or ''

                    resources = page.get('/Resources')
                    if isinstance(resources, IndirectObject):
                        resources = resources.get_object()

                    if not has_images and resources and '/XObject' in resources:
                        xObject = resources['/XObject']
                        if isinstance(xObject, IndirectObject):
                            xObject = xObject.get_object()
                        has_images = any(
                            xObject[obj].get('/Subtype') == '/Image'
                            for obj in xObject
                        )
                except Exception as e:
                    logger.warning(f"Error reading page {i} of {temp_file}: {e}")

        except Exception as e:
            logger.error(f"Failed to open PDF {file_path}: {e}")
            return ''
        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except:
                    pass

        # Fallback to OCR
        if not extracted_text.strip() and has_images:
            if self.ocr_handler:
                logger.info(f"OCR required for file: {file_path}")
                return self.ocr_handler.extract_text_from_pdf(file_path)
            else:
                logger.warning(f"The file {file_path} requires OCR but no handler was provided.")

        return extracted_text

    def _extract_pdf_with_metadata(self, file_path: str) -> Dict:
        """Extract PDF text with page-by-page metadata."""
        filename = Path(file_path).name
        pages_data = []
        temp_file = None

        try:
            temp_file = self._resave_pdf(file_path)
            reader = PdfReader(temp_file)

            for i, page in enumerate(reader.pages, 1):
                try:
                    page_text = page.extract_text() or ''
                    # Include all pages, even if empty
                    pages_data.append({"page": i, "content": page_text})
                except Exception as e:
                    logger.warning(f"Error reading page {i} of {temp_file}: {e}")
                    # Add empty page on error
                    pages_data.append({"page": i, "content": ""})

        except Exception as e:
            logger.error(f"Failed to open PDF {file_path}: {e}")
        finally:
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except:
                    pass

        # Fallback to OCR if no text found in any page
        if not any(page["content"].strip() for page in pages_data) and self.ocr_handler:
            logger.info(f"OCR required for file: {file_path}")
            ocr_pages = self.ocr_handler.extract_text_from_pdf_by_pages(file_path)
            if ocr_pages:
                pages_data = [{"page": i+1, "content": page_text} for i, page_text in enumerate(ocr_pages)]

        # Generate combined text and metadata dictionary
        all_text = " ".join(page["content"] for page in pages_data if page["content"].strip())

        if not all_text.strip():
            return {"text": "", "metadata": {}}

        # Create metadata dictionary structure
        metadata = {filename: {}}
        for page_data in pages_data:
            if page_data["content"].strip():  # Only include non-empty pages
                metadata[filename][page_data["page"]] = page_data["content"]

        return {
            "text": all_text,
            "metadata": metadata
        }

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        try:
            doc = docx.Document(file_path)
            return ' '.join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")
            return ""

    def _extract_txt(self, file_path: str) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error processing TXT file {file_path}: {e}")
                return ""

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        try:
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            return ""

    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML files."""
        encodings = ['utf-8', 'latin-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    soup = BeautifulSoup(file.read(), 'html.parser')
                    return soup.get_text(separator=' ', strip=True)
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.error(f"Error processing HTML file {file_path}: {e}")
                break
        return ""

    def _extract_odt(self, file_path: str) -> str:
        """Extract text from OpenDocument Text files."""
        try:
            doc = load(file_path)
            return ' '.join(
                teletype.extractText(element)
                for element in doc.getElementsByType(text.P)
            )
        except Exception as e:
            logger.error(f"Error processing ODT file {file_path}: {e}")
            return ""

    def _extract_rtf(self, file_path: str) -> str:
        """Extract text from RTF files."""
        try:
            with open(file_path, 'r') as file:
                return rtf_to_text(file.read(), errors='ignore')
        except Exception as e:
            logger.error(f"Error processing RTF file {file_path}: {e}")
            return ""

    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return ' '.join(
                    ' '.join(row)
                    for row in csv.reader(file)
                    if any(field.strip() for field in row)
                )
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {e}")
            return ""

    def _extract_xml(self, file_path: str) -> str:
        """Extract text from XML files."""
        try:
            tree = ET.parse(file_path)
            return ' '.join(
                elem.text.strip()
                for elem in tree.iter()
                if elem.text and elem.text.strip()
            )
        except Exception as e:
            logger.error(f"Error processing XML file {file_path}: {e}")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract evaluated text content from Excel files."""
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            texts.append(str(cell))
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing XLSX file {file_path}: {e}")
            return ""

    def _extract_xls(self, file_path: str) -> str:
        """Extract text from legacy Excel files."""
        try:
            book = xlrd.open_workbook(file_path, formatting_info=False)
            texts = []
            for sheet in book.sheets():
                for row_idx in range(sheet.nrows):
                    for cell in sheet.row(row_idx):
                        value = cell.value
                        if value and not str(value).startswith('='):
                            texts.append(str(value))
            return ' '.join(texts)
        except Exception as e:
            logger.error(f"Error processing XLS file {file_path}: {e}")
            return ""

    def _extract_ods(self, file_path: str) -> str:
        """Extract text from OpenDocument Spreadsheets."""
        try:
            doc = load(file_path)
            return '\n'.join(
                "".join(
                    child.data
                    for child in p.childNodes
                    if child.nodeType == child.TEXT_NODE
                )
                for p in doc.getElementsByType(P)
            )
        except Exception as e:
            logger.error(f"Error processing ODS file {file_path}: {e}")
            return ""

    def _extract_dbf(self, file_path: str) -> str:
        """Extract text from DBF database files."""
        try:
            dbf = DBF(file_path, load=True)
            return ' '.join(
                f"{key}: {value}"
                for record in dbf
                for key, value in record.items()
            )
        except Exception as e:
            logger.error(f"Error processing DBF file {file_path}: {e}")
            return ""
