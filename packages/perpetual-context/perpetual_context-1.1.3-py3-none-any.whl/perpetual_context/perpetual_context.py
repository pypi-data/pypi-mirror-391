# This code is an algorithm projected, architected and developed by Sapiens Technology®️ and aims to manage context window memory
# in Artificial Intelligence projects for language models. It manages context memory by saving and indexing the encoded dialogs
# for later consultation and return of excerpts referring to the input prompt, summarizing these excerpts when necessary to prevent
# the character sequence from exceeding the tokens limit previously established by the control variable.
# This makes it possible to establish a perpetual and limitless context even for language models with limited context windows.
class PerpetualContext:
    def __init__(self, display_error_point=True):
        try:
            self.__display_error_point = bool(display_error_point) if type(display_error_point) in (bool, int, float) else True
            def installModule(module_name='', version=None):
                try:
                    from subprocess import check_call, CalledProcessError
                    from sys import executable
                    module_name = str(module_name).strip()
                    module_name = f'{module_name}=={version}' if version else module_name
                    check_call([executable, '-m', 'pip', 'install', module_name])
                    print(f'Module {module_name} installed successfully!')
                except CalledProcessError as error:
                    print(f'ERROR installing the module "{module_name}": {error}')
                    print('Run the command:\npip install '+module_name)
            from traceback import print_exc
            from warnings import filterwarnings, catch_warnings, simplefilter
            from logging import getLogger, ERROR
            from unicodedata import normalize
            from re import sub, compile, search, split
            from io import BytesIO
            from os import path, makedirs, listdir, environ, unlink
            from shutil import rmtree
            from datetime import datetime
            from copy import deepcopy
            from tempfile import NamedTemporaryFile
            from base64 import b64encode, b64decode
            try: from count_tokens import count_tokens_in_string
            except:
                installModule(module_name='count-tokens', version='0.7.0')
                from count_tokens import count_tokens_in_string
            try: from requests import get
            except:
                installModule(module_name='requests', version='2.31.0')
                from requests import get
            try: from numpy import array
            except:
                installModule(module_name='numpy', version='1.25.2')
                from numpy import array
            try: from pandas import read_csv, read_excel
            except:
                installModule(module_name='pandas', version='2.2.2')
                from pandas import read_csv, read_excel
            try: from statistics import mean, median, stdev, variance
            except:
                installModule(module_name='statistics', version='1.0.3.5')
                from statistics import mean, median, stdev, variance
            try: from certifi import where
            except:
                installModule(module_name='certifi', version='2024.2.2')
                from certifi import where
            try: from tabulate import tabulate
            except:
                installModule(module_name='tabulate', version='0.9.0')
                from tabulate import tabulate
            try: from PyPDF2 import PdfReader
            except:
                installModule(module_name='PyPDF2', version='3.0.1')
                from PyPDF2 import PdfReader
            try: from fitz import open as fitz, Matrix
            except:
                installModule(module_name='PyMuPDF', version='1.24.5')
                from fitz import open as fitz, Matrix
            try: from docx import Document
            except:
                installModule(module_name='python-docx', version='1.1.0')
                from docx import Document
            try: from pptx import Presentation
            except:
                installModule(module_name='python-pptx', version='0.6.23')
                from pptx import Presentation
            try: from bs4 import BeautifulSoup
            except:
                installModule(module_name='beautifulsoup4', version='4.12.3')
                from bs4 import BeautifulSoup
            try: from youtubesearchpython import Video, ResultMode
            except:
                installModule(module_name='youtube-search-python', version='1.6.6')
                from youtubesearchpython import Video, ResultMode
            try: from youtube_transcript_api import YouTubeTranscriptApi
            except:
                installModule(module_name='youtube-transcript-api', version='0.6.2')
                from youtube_transcript_api import YouTubeTranscriptApi
            try: from PIL import Image
            except:
                installModule(module_name='pillow', version='10.3.0')
                from PIL import Image
            try: from easyocr import Reader
            except:
                installModule(module_name='easyocr', version='1.7.1')
                from easyocr import Reader
            try: from torch import cuda, no_grad
            except:
                installModule(module_name='torch', version='2.3.0')
                from torch import cuda, no_grad   
            try: from torchvision import models, transforms
            except:
                installModule(module_name='torchvision', version='0.18.0')
                from torchvision import models, transforms
            try: from webcolors import CSS3_HEX_TO_NAMES, hex_to_rgb, rgb_to_name
            except:
                installModule(module_name='webcolors', version='1.13')
                from webcolors import CSS3_HEX_TO_NAMES, hex_to_rgb, rgb_to_name
            try: from sklearn.cluster import KMeans
            except:
               installModule(module_name='scikit-learn', version='1.5.0')
               from sklearn.cluster import KMeans 
            try: from pydub import AudioSegment
            except:
                installModule(module_name='pydub', version='0.25.1')
                from pydub import AudioSegment
            try: from speech_recognition import Recognizer, AudioFile
            except:
                installModule(module_name='SpeechRecognition', version='3.10.3')
                from speech_recognition import Recognizer, AudioFile
            self.__count_tokens_in_string = count_tokens_in_string
            self.__get = get
            self.__array = array
            self.__read_csv, self.__read_excel = read_csv, read_excel
            self.__mean, self.__median, self.__stdev, self.__variance = mean, median, stdev, variance
            self.__where = where
            self.__tabulate = tabulate
            self.__PdfReader = PdfReader
            self.__fitz, self.__Matrix = fitz, Matrix
            self.__Document = Document
            self.__Presentation = Presentation
            self.__BeautifulSoup = BeautifulSoup
            self.__Video, self.__ResultMode = Video, ResultMode
            self.__YouTubeTranscriptApi = YouTubeTranscriptApi
            self.__Image = Image
            self.__Reader = Reader
            self.__cuda, self.__no_grad = cuda, no_grad
            self.__models, self.__transforms = models, transforms
            self.__CSS3_HEX_TO_NAMES, self.__hex_to_rgb, self.__rgb_to_name = CSS3_HEX_TO_NAMES, hex_to_rgb, rgb_to_name
            self.__KMeans = KMeans
            self.__AudioSegment = AudioSegment
            self.__Recognizer, self.__AudioFile = Recognizer, AudioFile
            self.__print_exc = print_exc
            self.__filterwarnings, self.__catch_warnings, self.__simplefilter = filterwarnings, catch_warnings, simplefilter
            self.__getLogger, self.__ERROR = getLogger, ERROR
            self.__normalize = normalize
            self.__sub, self.__compile, self.__search, self.__split = sub, compile, search, split
            self.__BytesIO = BytesIO
            self.__path, self.__makedirs, self.__listdir, self.__environ, self.__unlink = path, makedirs, listdir, environ, unlink
            self.__rmtree = rmtree
            self.__datetime = datetime
            self.__deepcopy = deepcopy
            self.__NamedTemporaryFile = NamedTemporaryFile
            self.__b64encode, self.__b64decode = b64encode, b64decode 
            languages1 = ['en', 'pt', 'es', 'fr', 'de', 'it', 'ru', 'ar', 'hi', 'zh-Hant', 'zh-Hans', 'af', 'ak', 'sq', 'am', 'hy', 'as', 'ay', 'az', 'bn', 'eu', 'be', 'bho', 'bs', 'bg']
            languages2 = ['my', 'ca', 'ceb', 'co', 'hr', 'cs', 'da', 'dv', 'nl', 'eo', 'et', 'ee', 'fil', 'fi', 'gl', 'lg', 'ka', 'el', 'gn', 'gu', 'ht', 'ha', 'haw', 'iw', 'hmn']
            languages3 = ['hu', 'is', 'ig', 'id', 'ga', 'ja', 'jv', 'kn', 'kk', 'km', 'rw', 'ko', 'kri', 'ku', 'ky', 'lo', 'la', 'lv', 'ln', 'lt', 'lb', 'mk', 'mg', 'ms', 'ml']
            languages4 = ['mt', 'mi', 'mr', 'mn', 'ne', 'nso', 'no', 'ny', 'or', 'om', 'ps', 'fa', 'pl', 'pa', 'qu', 'ro', 'sm', 'sa', 'gd', 'sr', 'sn', 'sd', 'si', 'sk', 'sl']
            languages5 = ['so', 'st', 'su', 'sw', 'sv', 'tg', 'ta', 'tt', 'te', 'th', 'ti', 'ts', 'tr', 'tk', 'uk', 'ur', 'ug', 'uz', 'vi', 'cy', 'fy', 'xh', 'yi', 'yo', 'zu']
            self.__languages = languages1+languages2+languages3+languages4+languages5
            self.__object_names = ['__background__', 'person', 'bicycle', 'car', 'motorcycle', 'airplane', 'bus', 'train', 'truck', 'boat', 'traffic light', 'fire hydrant', 'N/A', 'stop sign',
            'parking meter', 'bench', 'bird', 'cat', 'dog', 'horse', 'sheep', 'cow', 'elephant', 'bear', 'zebra', 'giraffe', 'N/A', 'backpack', 'umbrella', 'N/A', 'N/A', 'handbag',
            'tie', 'suitcase', 'frisbee', 'skis', 'snowboard', 'sports ball', 'kite', 'baseball bat', 'baseball glove', 'skateboard', 'surfboard', 'tennis racket', 'bottle', 'N/A', 'wine glass',
            'cup', 'fork', 'knife', 'spoon', 'bowl', 'banana', 'apple', 'sandwich', 'orange', 'broccoli', 'carrot', 'hot dog', 'pizza', 'donut', 'cake', 'chair', 'couch', 'potted plant', 'bed',
            'N/A', 'dining table', 'N/A', 'N/A', 'toilet', 'N/A', 'tv', 'laptop', 'mouse', 'remote', 'keyboard', 'cell phone', 'microwave', 'oven', 'toaster', 'sink', 'refrigerator', 'N/A',
            'book', 'clock', 'vase', 'scissors', 'teddy bear', 'hair drier', 'toothbrush']
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__init__: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
    def __getRootDirectory(self):
        try: return str(self.__path.dirname(self.__path.realpath(__file__)).replace('\\', '/')+'/').strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getRootDirectory: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return './'
    def __utf8ToBase64(self, utf8_string=''):
        try:
            utf8_string = str(utf8_string)
            utf8_bytes = utf8_string.encode('utf-8')
            base64_bytes = self.__b64encode(utf8_bytes)
            base64_string = base64_bytes.decode('utf-8')
            return base64_string
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__utf8ToBase64: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return utf8_string
    def __base64ToUtf8(self, base64_string=''):
        try:
            base64_string = str(base64_string)
            base64_bytes = base64_string.encode('utf-8')
            utf8_bytes = self.__b64decode(base64_bytes)
            utf8_string = utf8_bytes.decode('utf-8')
            return utf8_string
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__base64ToUtf8: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return base64_string
    def __countTokens(self, string=''):
        try: return self.__count_tokens_in_string(str(string).replace('<|', '<!').replace('|>', '!>')) if '|' in str(string) else self.__count_tokens_in_string(str(string))
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__countTokens: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return len(str(string))
    def __removeAccentsAndSpecialChars(self, string=''):
        try: return self.__sub(r'[^\w\s]', '', self.__normalize('NFKD', str(string).strip()).encode('ASCII', 'ignore').decode('utf-8'))
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__removeAccentsAndSpecialChars: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return string
    def __isNumber(self, string=''):
        try: return type(float(str(string).strip())) == float
        except: return False
    def __getKeyWords(self, string=''):
        try:
            key_words = []
            string = str(string).strip()
            if len(string) > 0:
                def isAcronym(string=''): return str(string).strip().isupper()
                def isImportant(string=''): return len(str(string).strip()) > 4
                def isKeyWord(string=''): return self.__isNumber(string) or isAcronym(string) or isImportant(string)
                for word in string.split():
                    if isKeyWord(word): key_words.append(self.__removeAccentsAndSpecialChars(string=word).lower().strip())
            return key_words
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getKeyWords: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return []
    def __getAllTokens(self, user_id=0, dialog_id=0):
        try:
            all_tokens = 0
            user_id, dialog_id = str(user_id).strip(), str(dialog_id).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            user_context = f'{context_directory}{user_id}/'
            tokens_path, file_text = f'{user_context}{dialog_id}.tokens', '0'
            if self.__path.exists(tokens_path):
                with open(tokens_path, encoding='utf-8', errors='ignore') as file:
                    file_text = str(file.read()).strip()
                    file.close()
                all_tokens = max((0, int(round(float(file_text)))))
            return all_tokens
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getAllTokens: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return 0
    def __saveAllTokens(self, user_id=0, dialog_id=0, all_tokens=0):
        try:
            user_id, dialog_id = str(user_id).strip(), str(dialog_id).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            all_tokens = max((0, int(round(float(all_tokens))))) if type(all_tokens) in (int, float) else 0
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            user_context = f'{context_directory}{user_id}/'
            tokens_path = f'{user_context}{dialog_id}.tokens'
            write = open(tokens_path, 'w', encoding='utf-8', errors='ignore')
            write.write(str(all_tokens))
            write.close()
            return True
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__saveAllTokens: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return False
    def __getIndex(self, user_id=0, dialog_id=0, index=0):
        try:
            conversation = '\n<|conversation|>\n'
            user_id, dialog_id = str(user_id).strip(), str(dialog_id).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            index = max((0, int(round(float(index))))) if type(index) in (int, float) else 0
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            user_context = f'{context_directory}{user_id}/'
            dialogue_context = f'{user_context}{dialog_id}/'
            index_path = f'{dialogue_context}{index}.index'
            if self.__path.isdir(dialogue_context) and self.__path.exists(index_path):
                with open(index_path, encoding='utf-8', errors='ignore') as file:
                    conversation = self.__base64ToUtf8(base64_string=str(file.read()).strip())
                    file.close()
            return conversation
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getIndex: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return '\n<|conversation|>\n'
    def __getRelatedContext(self, user_id=0, dialog_id=0, key_words=[], max_tokens=1, indexes=[]):
        try:
            context_window = []
            user_id, dialog_id = str(user_id).strip(), str(dialog_id).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            key_words = list(key_words) if type(key_words) in (tuple, list) else []
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            indexes = list(indexes) if type(indexes) in (tuple, list) else []
            half_of_the_tokens = int((max_tokens/3)/2)
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            user_context = f'{context_directory}{user_id}/'
            dialogue_context = f'{user_context}{dialog_id}/'
            dialog_length = len(self.__listdir(dialogue_context))
            if len(indexes) > 0:
                for index in indexes:
                    if index < 0: index = dialog_length-abs(index)
                    if self.__path.exists(f'{dialogue_context}{index}.index'): context_window.append(self.__getIndex(user_id=user_id, dialog_id=dialog_id, index=index))
            if len(context_window) > 0: return context_window
            if self.__path.exists(f'{dialogue_context}0.index'): context_window.append(self.__getIndex(user_id=user_id, dialog_id=dialog_id, index=0))
            default_index = max((0, int(round(float(dialog_length/2)))-1))
            def summarizeInputAndOutput(conversation=''):
                tokens_number = self.__countTokens(string=conversation.replace('<|conversation|>', ' '))
                if tokens_number > max_tokens:
                    _input, _output = conversation.split('<|conversation|>')
                    _input, _output = self.__getSummaryText(text=_input, max_tokens=half_of_the_tokens), self.__getSummaryText(text=_output, max_tokens=half_of_the_tokens)
                    conversation = _input+'\n<|conversation|>\n'+_output
                return conversation.strip()
            if len(key_words) > 0:
                maximum_count, maximum_index = 0, default_index
                for index in range(1, dialog_length-1):
                    if self.__path.exists(f'{dialogue_context}{index}.index'):
                        context_tokens = self.__removeAccentsAndSpecialChars(self.__getIndex(user_id=user_id, dialog_id=dialog_id, index=index).replace('<|conversation|>', ' ')).strip().split()
                        context_tokens, token_count = [token.lower().strip() for token in context_tokens if len(token.strip()) > 0], 0
                        for key_word in key_words:
                            if key_word in context_tokens: token_count += 1
                        if token_count > maximum_count: maximum_count, maximum_index = token_count, index
                if self.__path.exists(f'{dialogue_context}{maximum_index}.index'):
                    indexed_content = self.__getIndex(user_id=user_id, dialog_id=dialog_id, index=maximum_index)
                    if indexed_content not in context_window: context_window.append(summarizeInputAndOutput(indexed_content))
            elif self.__path.exists(f'{dialogue_context}{default_index}.index'):
                indexed_content = self.__getIndex(user_id=user_id, dialog_id=dialog_id, index=default_index)
                if indexed_content not in context_window: context_window.append(summarizeInputAndOutput(indexed_content))
            if self.__path.exists(f'{dialogue_context}{dialog_length-1}.index'):
                indexed_content = self.__getIndex(user_id=user_id, dialog_id=dialog_id, index=dialog_length-1)
                if indexed_content not in context_window: context_window.append(summarizeInputAndOutput(indexed_content))
            return context_window
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getRelatedContext: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return []
    def __getNumberOfDialogues(self, user_id=0, dialog_id=0):
        try:
            user_id, dialog_id = str(user_id).strip(), str(dialog_id).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            user_context = f'{context_directory}{user_id}/'
            dialogue_context = f'{user_context}{dialog_id}/'
            return len(self.__listdir(dialogue_context)) if self.__path.isdir(dialogue_context) else 0
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getNumberOfDialogues: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return 0
    def __isTitle(self, string=''):
        try:
            is_title = False
            string = str(string).strip()
            length_of_text_snippet = len(string)
            if length_of_text_snippet <= 0: return is_title
            first_character_of_the_string = string[0]
            last_character_of_the_string = string[-1]
            finishers_not_considered = (',', ';', '.')
            title_launcher = first_character_of_the_string.isupper() or not first_character_of_the_string.isalpha()
            without_finisher = last_character_of_the_string not in finishers_not_considered
            if length_of_text_snippet <= 23 and title_launcher and without_finisher: return True
            def notIsPhrase(): return length_of_text_snippet <= 162 and title_launcher and without_finisher
            if notIsPhrase():
                if ' ' in string:
                    tokens = string.split()
                    is_title = sum([1 if (len(token.strip()) > 0 and (token.strip()[0].isupper() or not token.strip()[0].isalpha())) or len(token.strip()) <= 4 or self.__isNumber(token) else 0 for token in tokens]) == len(tokens)
                elif string.isupper(): is_title = True
            return is_title
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__isTitle: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return False
    def __getBase64ImageFromPDF(self, file_path='', page_index=0):
        try:
            image_dictionary = {'base64': '', 'extension': '.png'}
            file_path = str(file_path).strip()
            page_index = max((0, int(page_index))) if type(page_index) in (int, float) else 0
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                response = self.__get(file_path)
                pdf_bytes = self.__BytesIO(response.content)
                document = self.__fitz(stream=pdf_bytes, filetype='pdf')
            elif not self.__path.exists(file_path):
                print(f'The path to the "{file_path}" file does not exist.')
                return image_dictionary
            else: document = self.__fitz(file_path)
            document_page = document[page_index]
            zoom_x, zoom_y = 2, 2
            matrix = self.__Matrix(zoom_x, zoom_y)
            page_image = document_page.get_pixmap(matrix=matrix)
            image_bytes = page_image.tobytes('png')
            image_base64 = self.__b64encode(image_bytes).decode('utf-8')
            image_dictionary['base64'] = str(image_base64)
            return image_dictionary
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getBase64ImageFromPDF: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return {'base64': '', 'extension': '.png'}
    def __getColumnType(self, vector=[]):
        _, result_type = None, ''
        irrelevants = ('', 'none', 'null', 'undefined', 'nan')
        for element in vector:
            element = str(element).strip()
            try:
                if element.lower() in ('false', 'true'): result_type = 'boolean'
                elif '.' in element: _, result_type = str(float(element)), 'float number'
                elif element.isdigit(): _, result_type = str(int(element)), 'integer number'
                elif element not in irrelevants: _, result_type = str(complex(element)), 'complex number'
                else: _, result_type = str(_), 'text'
            except:  _, result_type = element, 'text'
            if len(element) > 0 and _.lower() not in irrelevants: break
        if result_type == 'text' and len(_) == 1: result_type = 'character'
        return result_type if len(result_type) > 0 else 'no type'
    def __getOrderingOfValues(self, vector=[]):
        if len(vector) < 1: return []
        vector = [str(x) if type(x) not in (bool, int, float, str) else x for x in vector]
        return sorted(vector, key=lambda x: (isinstance(x, str), x))
    def __getRepetitionOfValues(self, vector=[], function=min):
        if len(vector) < 1: return 'does not repeat'
        vector = self.__getOrderingOfValues(vector=vector)
        counts = [vector.count(x) for x in vector]
        repetition = vector[counts.index(function(counts))]
        return 'does not repeat' if sum(counts) == len(counts) else repetition
    def __getExtremityOfValues(self, vector=[]):
        if len(vector) < 1: return (0, 0)
        vector = vector = self.__getOrderingOfValues(vector=vector)
        return (vector[0], vector[-1])
    def __getAverageValue(self, vector=[]):
        if len(vector) < 1: return 0
        vector = self.__getOrderingOfValues(vector=vector)
        vector_length = len(vector)
        average_index = min((max((0, int(round((vector_length-1)/2)))), vector_length))
        return vector[average_index]
    def __getAverageOfValues(self, vector=[]):
        try: return self.__mean(vector) if len(vector) > 0 else 0
        except: return self.__getAverageValue(vector=vector)
    def __getMedianOfValues(self, vector=[]):
        try: return self.__median(vector) if len(vector) > 0 else 0
        except: return self.__getAverageValue(vector=vector)
    def __getVectorWithUniqueElements(self, vector=[]):
        if len(vector) < 1: return []
        vector = [str(x) if type(x) not in (type(None), int, float, str) else x for x in vector]
        temporary_vector = []
        for element in vector:
            if element not in temporary_vector: temporary_vector.append(element)
        return temporary_vector
    def __getStandardDeviationOfValues(self, vector=[]):
        try: return self.__stdev(vector) if len(vector) > 1 else 0
        except:
            vector = [str(x) for x in vector]
            unique_elements = self.__getVectorWithUniqueElements(vector=vector)
            for index, element in enumerate(unique_elements): vector = [index if x == element else x for x in vector]
            return self.__stdev(vector)
    def __getVarianceOfValues(self, vector=[]):
        try: return self.__variance(vector) if len(vector) > 1 else 0
        except:
            vector = [str(x) for x in vector]
            unique_elements = self.__getVectorWithUniqueElements(vector=vector)
            for index, element in enumerate(unique_elements): vector = [index if x == element else x for x in vector]
            return self.__variance(vector)
    def __getTableSummary(self, _table=[], _column_metrics=[], _metrics_indexes=()):
        result_table = []
        if len(_metrics_indexes) > 0:
            for index_x, column_metric in enumerate(_column_metrics):
                temporary_metric = []
                for index_y, metric in enumerate(column_metric):
                    if index_y in _metrics_indexes: temporary_metric.append(metric)
                _column_metrics[index_x] = tuple(temporary_metric)
        for column_index, column_metric in enumerate(_column_metrics):
            for row in _table:
                column = list(zip(*result_table))[column_index] if len(result_table) > 0 else []
                if row[column_index] in column_metric and row[column_index] not in column and row not in result_table: result_table.append(row)
        return result_table
    def __getSummaryText(self, text='', max_tokens=1):
        try:
            summary_result = ''
            text = str(text).replace('  ', ' ').strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            tokens_number = self.__countTokens(string=text)
            if tokens_number <= max_tokens: return text
            title_separators, default_title_separator, possible_titles, titles = ('\n', '?', '!', ':'), '', [], []
            for title_separator in title_separators:
                if title_separator in text:
                    possible_titles = text.split(title_separator)
                    default_title_separator = title_separator
                    break
            for possible_title in possible_titles:
                if self.__isTitle(possible_title): titles.append(possible_title)
            separators, default_separator = ('.\n\n', '\n\n', '.\n', '\n', '. ', '.', ';', '?', '!', ':', ','), ' '
            def getBeginningMiddleAndEnd(text=''):
                beginning_middle_and_end = text = str(text).strip()
                chosen_separator = ' '
                for separator in separators:
                    if text.count(separator) > 3:
                        chosen_separator = separator
                        break
                vectorized_text = text.split(chosen_separator)
                if len(vectorized_text) > 3:
                    middle_index, temporary_vector = min((len(vectorized_text)-1, max((0, int(round(float(len(vectorized_text)/2)))-1)))), []
                    if vectorized_text[0] not in temporary_vector: temporary_vector.append(vectorized_text[0])
                    if vectorized_text[middle_index] not in temporary_vector: temporary_vector.append(vectorized_text[middle_index])
                    if vectorized_text[-1] not in temporary_vector: temporary_vector.append(vectorized_text[-1])
                    vectorized_text = temporary_vector.copy()
                    beginning_middle_and_end = chosen_separator.join(vectorized_text)
                return beginning_middle_and_end
            if len(titles) > 0:
                title_contents, initialization = [], '<|initialization|>'
                for title in titles: text = text.replace(title, initialization, 1)
                before_the_first_title, text_tokens = '', text.split(initialization) if initialization in text else []
                if not text.startswith(initialization):
                    before_the_first_title = text_tokens[0]
                    title_contents = text_tokens[1:]
                else:
                    if len(text_tokens) > 0 and len(text_tokens[0].strip()) <= 0: text_tokens = text_tokens[1:]
                    title_contents = text_tokens.copy()
                if len(titles) == len(title_contents):
                    for title_index, title_content in enumerate(title_contents):
                        for separator in separators:
                            if title_content.count(separator) > 3:
                                default_separator = separator
                                break
                        part_summary = title_content
                        content_parts = title_content.split(default_separator)
                        if len(content_parts) > 3:
                            middle_index, temporary_vector = min((len(content_parts)-1, max((0, int(round(float(len(content_parts)/2)))-1)))), []
                            if content_parts[0] not in temporary_vector: temporary_vector.append(content_parts[0])
                            if content_parts[middle_index] not in temporary_vector: temporary_vector.append(content_parts[middle_index])
                            if content_parts[-1] not in temporary_vector: temporary_vector.append(content_parts[-1])
                            content_parts = temporary_vector.copy()
                            part_summary = default_separator.join(content_parts)
                        summary_result += str(titles[title_index]+default_title_separator+part_summary).strip()+'\n\n'
                    if len(before_the_first_title.strip()) > 0: summary_result = before_the_first_title.strip()+'\n\n'+summary_result
                else:
                    for title in titles: text = text.replace('<|initialization|>', title, 1)
                    summary_result = getBeginningMiddleAndEnd(text=text)
            else: summary_result = getBeginningMiddleAndEnd(text=text)
            return summary_result.strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getSummaryText: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return text
    def __getBeginningAndEnd(self, string='', max_tokens=1, separator=''):
        try:
            string = str(string).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            if type(separator) != str: separator = ''
            half = int(max_tokens/2)
            return string[:half]+separator+string[::-1][:half][::-1]
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getBeginningAndEnd: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return string
    def __getBeginningMiddleAndEnd(self, string='', max_tokens=1, separator=''):
        try:
            string = str(string).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            if type(separator) != str: separator = ''
            parts_length = max_tokens//3
            beginning = string[:parts_length]
            middle = string[len(string)//2-parts_length//2:len(string)//2+parts_length//2]
            end = string[-parts_length:]
            return beginning+separator+middle+separator+end
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__getBeginningMiddleAndEnd: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return self.__getBeginningAndEnd(string=string, max_tokens=max_tokens, separator=separator)
    def __formatPageText(self, page_text=''):
        try:
            page_text = page_text.strip()
            page_text = page_text.replace('\n\n\n\n\n', '\n').replace('\n\n\n\n', '\n').replace('\n\n\n', '\n').replace('\n\n', '\n')
            page_text = page_text.replace('-\n', '').replace('- \n', '').replace('     ', ' ').replace('    ', ' ').replace('   ', ' ').replace('  ', ' ')
            return page_text.strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__formatPageText: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return page_text
    def __removeLinkDelimiters(self, string=''): return self.__compile(r'^[\'\"]|[\'\"]$').sub('', self.__sub(r'[,.;:!?]*(?:\s*\.\.\.)?$', '', self.__sub(r"^[\'\"](.*?)[\'\"]$", r"\1", str(string).strip())))
    def __itsCode(self, text=''):
        try:
            def containsCode(input_string):
                python_patterns = [r'\bdef\b\s+\w+\s*\(.*\):', r'\bclass\b\s+\w+\s*\(.*\):', r'print\s*\(.*\)', r'\bfor\b\s+\w+\s+in\s+.*:', r'\bwhile\b\s+.*:']
                javascript_patterns = [r'\bfunction\b\s+\w+\s*\(.*\)\s*\{', r'console\.log\s*\(.*\);', r'var\s+\w+\s*=\s*.*;', r'\bfor\b\s*\(.*\)\s*\{', r'\bwhile\b\s*\(.*\)\s*\{']
                html_patterns = [r'<\w+\s*.*>', r'</\w+>', r'<\w+\s*.*>\s*.*\s*</\w+>']
                c_patterns = [r'#include\s*<.*>', r'\bint\b\s+\w+\s*\(.*\)\s*\{', r'\bfor\b\s*\(.*\)\s*\{', r'\bwhile\b\s*\(.*\)\s*\{']
                cpp_patterns = [r'#include\s*<.*>', r'\bclass\b\s+\w+\s*\{', r'\bint\b\s+\w+\s*\(.*\)\s*\{', r'\bfor\b\s*\(.*\)\s*\{', r'\bwhile\b\s*\(.*\)\s*\{']
                java_patterns = [r'\bclass\b\s+\w+\s*\{', r'\bpublic\b\s+\bstatic\b\s+\bvoid\b\s+\w+\s*\(.*\)\s*\{', r'System\.out\.println\s*\(.*\);', r'\bfor\b\s*\(.*\)\s*\{', r'\bwhile\b\s*\(.*\)\s*\{']
                csharp_patterns = [r'\bclass\b\s+\w+\s*\{', r'\bpublic\b\s+\bstatic\b\s+\bvoid\b\s+\w+\s*\(.*\)\s*\{', r'Console\.WriteLine\s*\(.*\);', r'\bfor\b\s*\(.*\)\s*\{', r'\bwhile\b\s*\(.*\)\s*\{']
                php_patterns = [r'<\?php', r'\bfunction\b\s+\w+\s*\(.*\)\s*\{', r'echo\s+.*;', r'\bfor\b\s*\(.*\)\s*\{', r'\bwhile\b\s*\(.*\)\s*\{']
                all_patterns = python_patterns+javascript_patterns+html_patterns+c_patterns+cpp_patterns+java_patterns+csharp_patterns+php_patterns
                for pattern in all_patterns:
                    if self.__search(pattern, input_string): return True
                return False
            return containsCode(input_string=str(text).strip())
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.__itsCode: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return False
    def getFileType(self, file_path=''): return self.__path.splitext(str(file_path).strip())[1][1:].strip()
    def countTokens(self, string=''): return self.__countTokens(string=string)
    def getKeyWords(self, string=''): return self.__getKeyWords(string=string)
    def getBeginningAndEnd(self, string='', max_tokens=1, separator=''): return self.__getBeginningAndEnd(string=string, max_tokens=max_tokens, separator=separator)
    def getBeginningMiddleAndEnd(self, string='', max_tokens=1, separator=''): return self.__getBeginningMiddleAndEnd(string=string, max_tokens=max_tokens, separator=separator)
    def getSummaryCode(self, text='', max_tokens=1):
        try:
            text = str(text).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            summary_result, found_code = text, False
            if '\n' in summary_result:
                lines_of_code, summary_lines = summary_result.split('\n'), []
                encodings = ('import', 'include', 'require', 'using', 'class', 'func', 'function', 'def', 'procedure', '{', '}', 'begin', 'end', 'then', 'new', 'if', 'elif', 'else', 'else:', 'else{', 'switch', 'case', 'for',
                             'while', 'do', 'loop', 'lambda', 'map', 'redux', '<!doctype', '<html', '<head', '<meta', '<title', '<link', '<script', '<body', '<div', '//', '/*', '*/', '#', '"""', "'''", '<!--', '-->')
                multiplier, initial_character = 0, ''
                for index, line_of_code in enumerate(lines_of_code):
                    if index > 0:
                        tokens = line_of_code.split()
                        multiplier = len(line_of_code) - len(line_of_code.lstrip(initial_character))
                        for token in tokens:
                            if token.lower().strip() in encodings:
                                summary_lines.append(line_of_code)
                                found_code = True
                            elif len(summary_lines) > 0 and '...' not in summary_lines[-1]: summary_lines.append((initial_character*multiplier)+'...')
                        if line_of_code.startswith(' '): initial_character = ' '
                        elif line_of_code.startswith('\t'): initial_character = '\t'
                summary_result = '\n'.join(summary_lines)
            tokens_number = self.__countTokens(string=summary_result)
            if tokens_number > max_tokens or not found_code: summary_result = self.__getBeginningMiddleAndEnd(string=text, max_tokens=max_tokens, separator='\n...\n')
            tokens_number = self.__countTokens(string=summary_result)
            if tokens_number > max_tokens: summary_result = self.__getBeginningMiddleAndEnd(string=text, max_tokens=max_tokens)
            return summary_result.strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryCode: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return self.__getBeginningMiddleAndEnd(string=text, max_tokens=max_tokens)
    def getSummaryText(self, text='', max_tokens=1):
        try:
            if self.__itsCode(text=text): return self.getSummaryCode(text=text, max_tokens=max_tokens)
            text = str(text).replace('    ', ' ').replace('   ', ' ').replace('  ', ' ').strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            summary_result = self.__getSummaryText(text=text, max_tokens=max_tokens)
            if len(text) <= len(summary_result): return self.__getBeginningMiddleAndEnd(string=text, max_tokens=max_tokens)
            tokens_number, attempts = self.__countTokens(string=summary_result), 0
            limit = len(summary_result.split())
            while tokens_number > max_tokens:
                summary_result = self.__getSummaryText(text=summary_result, max_tokens=max_tokens)
                tokens_number = self.__countTokens(string=summary_result)
                attempts += 1
                if attempts > limit: break
            tokens_number = self.__countTokens(string=summary_result)
            if tokens_number > max_tokens: summary_result = self.__getBeginningMiddleAndEnd(string=summary_result, max_tokens=max_tokens)
            return summary_result.strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryText: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return self.__getBeginningMiddleAndEnd(string=text, max_tokens=max_tokens)
    def getSummaryTXT(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path, document_text = str(file_path).strip(), ''
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                document_text = str(self.__get(file_path).text).strip()
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                with open(file_path, 'r') as file: document_text = str(file.read()).strip()
            if self.__itsCode(text=document_text):
                file_type = self.getFileType(file_path=file_path).upper()
                file_type_description = f'*{file_type} File*\n\n'
                document_text = self.getSummaryCode(text=document_text, max_tokens=max_tokens)
            else:
                file_type_description = '*TXT File*\n\n'
                document_text = self.getSummaryText(text=document_text, max_tokens=max_tokens)
            tokens_number = self.__countTokens(string=document_text)
            if tokens_number > max_tokens: document_text = self.__getBeginningMiddleAndEnd(string=document_text, max_tokens=max_tokens)
            result_text = file_type_description+document_text.strip()
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryTXT: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def imageToBase64(self, file_path=''):
        try:
            encoding_result = {'base64_string': '', 'image_type': 'jpg'}
            file_path = str(file_path).strip()
            image_type, content = 'jpg', None
            image_type = self.getFileType(file_path=file_path)
            if len(image_type) > 0: encoding_result['image_type'] = image_type
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                content = self.__get(file_path, stream=False).content
            if content != None: image_data = content
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return encoding_result
                with open(file_path, 'rb') as image_file: image_data = image_file.read()
            base64_encoded = self.__b64encode(image_data)
            base64_string = base64_encoded.decode('utf-8')
            encoding_result['base64_string'] = base64_string
            return encoding_result
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.imageToBase64: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return {'base64_string': '', 'image_type': ''}
    def saveBase64Image(self, base64_string='', file_path='', image_name='', extension=''):
        try:
            base64_string, file_path = str(base64_string).strip(), str(file_path).strip()
            if len(file_path) > 0 and file_path[-1] != '/': file_path = file_path+'/'
            image_name, extension = str(image_name).strip(), str(extension).strip()
            def getName(): return str(self.__datetime.now().timestamp()).replace('.', '').strip()
            if len(image_name) <= 0: image_name = getName()
            if len(extension) <= 1: extension = '.png' 
            if extension[0] != '.': extension = '.'+str(extension)
            if len(base64_string) <= 0: return False
            image_data = self.__b64decode(base64_string)
            file_path = file_path+image_name+extension
            with open(file_path, 'wb') as file: file.write(image_data)
            return self.__path.exists(file_path)
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.saveBase64Image: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return False
    def saveImageFromPDF(self, pdf_path='', image_path='', image_name='', extension='', page_index=0):
        try:
            pdf_path, image_path, image_name, extension = str(pdf_path).strip(), str(image_path).strip(), str(image_name).strip(), str(extension).strip()
            page_index = max((0, int(page_index))) if type(page_index) in (int, float) else 0
            image_dictionary = self.__getBase64ImageFromPDF(file_path=pdf_path, page_index=page_index)
            base64_string = str(image_dictionary['base64'])
            if len(extension) <= 0: extension = str(image_dictionary['extension']).strip()
            return self.saveBase64Image(base64_string=base64_string, file_path=image_path, image_name=image_name, extension=extension)
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.saveImageFromPDF: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return False
    def countPDFPages(self, file_path=''):
        try:
            pages_number = -1
            file_path = str(file_path).strip()
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                pdf_reader = self.__PdfReader(self.__BytesIO(self.__get(file_path).content))
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                pdf_file = open(file_path, 'rb')
                pdf_reader = self.__PdfReader(pdf_file)
            pages_number = len(pdf_reader.pages)
            return pages_number
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.countPDFPages: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return -1
    def getSummaryYouTube(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path = str(file_path).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            file_path = self.__removeLinkDelimiters(string=file_path)
            match = self.__search(r"(?:\?v=|\.be\/|\/shorts\/)([^&?/]+)", file_path)
            if match: file_path = match.group(1)
            file_path, description = self.__sub(r'\?feature=share', '', file_path), ''
            try:
                try: video = self.__Video.getInfo(file_path, mode=self.__ResultMode.json)
                except:
                    file_path = self.__split(r'\?', file_path)[0].strip() if self.__search(r'\?', file_path) else file_path.strip()
                    video = self.__Video.getInfo(file_path, mode=self.__ResultMode.json)
                channel_name = str(video["channel"]["name"]).strip() if ('channel' in video and 'name' in video['channel']) else ''
                title = str(video["title"]).strip() if 'title' in video else ''
                publish_date = str(video["publishDate"]).strip() if 'publishDate' in video else ''
                description = str(video["description"]).strip() if 'description' in video else ''
                if len(channel_name) > 0: result_text += f'**YouTube Channel Name:** {channel_name}\n'
                if len(title) > 0: result_text += f'**Video Title:** {title}\n'
                if len(publish_date) > 0: result_text += f'**Video Date:** {publish_date}\n'
                if len(description) > 0: result_text += f'**Video Description:**\n{description}\n'
            except: result_text = ''
            transcript, transcription = self.__YouTubeTranscriptApi.get_transcript(file_path, languages=self.__languages), ''
            for speech in transcript: transcription += f'**SPEECH:** {speech["text"]} - **START:** {speech["start"]} - **DURATION:** {speech["duration"]}\n'
            transcription = transcription.strip()
            def getSummaryValue(string='', max_tokens=1):
                separator = '\n'
                for divider in range(5, 0, -1):
                    tokens_number = self.__countTokens(string=string)
                    if tokens_number > max_tokens: string = self.__getBeginningMiddleAndEnd(string=string, max_tokens=max_tokens//divider, separator=separator)
                    else: break
                rows, temporary_rows = string.split(separator), []
                for row in rows:
                    if row.count('**SPEECH:**') == 1 and row.count('**START:**') == 1 and row.count('**DURATION:**') == 1: temporary_rows.append(row)
                if len(temporary_rows) > 0: string = separator.join(temporary_rows)
                return string
            transcription = getSummaryValue(string=transcription, max_tokens=max_tokens)
            if len(transcription) > 0: result_text += f'\n## CONTENT OF WHAT IS SAID IN THE VIDEO\n{transcription}'
            tokens_number = self.__countTokens(string=result_text)
            if tokens_number > max_tokens:
                result_text = result_text.replace(f'**Video Description:**\n{description}\n', '')
                tokens_number = self.__countTokens(string=result_text)
                if tokens_number > max_tokens: result_text = self.__getBeginningMiddleAndEnd(string=result_text, max_tokens=max_tokens)
            return result_text.strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryYouTube: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryWEBPage(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path = str(file_path).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            file_path = self.__removeLinkDelimiters(string=file_path)
            first_text, second_text = '', ''
            response = self.__get(file_path)
            if response.status_code == 200:
                source_code = self.__BeautifulSoup(response.content, 'html.parser')
                paragraphs = source_code.find_all('p')
                for paragraph in paragraphs: first_text += paragraph.get_text().strip()+'\n'
                first_text = str(first_text).strip()
                headings = source_code.find_all(self.__compile('^h[1-6]$'))
                for heading in headings:
                    second_text += '\n### '+heading.get_text().strip()+'\n'
                    paragraphs = heading.find_next_siblings('p')
                    for paragraph in paragraphs: second_text += paragraph.get_text().strip()+'\n'
                second_text = str(second_text).strip()
            result_text = first_text if len(first_text) > len(second_text) else second_text
            if len(result_text) > 0: result_text = f'## Contents of the WEB address: {file_path}\n\n{result_text}'
            else: result_text = file_path
            tokens_number = self.__countTokens(string=result_text)
            if tokens_number > max_tokens: result_text = self.getSummaryText(text=result_text, max_tokens=max_tokens)
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryWEBPage: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryPDF(self, file_path='', max_tokens=1, main_page=None):
        try:
            result_text = ''
            file_path, document_text = str(file_path).strip(), ''
            original_max_tokens = max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            main_page = max(1, int(main_page)) if type(main_page) in (int, float) else None
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                pdf_reader = self.__PdfReader(self.__BytesIO(self.__get(file_path).content))
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                pdf_file = open(file_path, 'rb')
                pdf_reader = self.__PdfReader(pdf_file)
            pages_number = len(pdf_reader.pages)
            file_type_description = '*PDF File*\n\n'
            if main_page != None:
                if main_page > pages_number: main_page = pages_number
                main_page_text = str(pdf_reader.pages[main_page-1].extract_text()).strip()
                tokens_number = self.__countTokens(string=main_page_text)
                if tokens_number < max_tokens: max_tokens = max_tokens-tokens_number
            max_tokens = int(((max_tokens-len(file_type_description))/pages_number)-len(f'PAGE {pages_number}'))
            if original_max_tokens > pages_number and max_tokens > 0:
                for page_index in range(pages_number):
                    try:
                        page_reader = pdf_reader.pages[page_index]
                        page_text = str(page_reader.extract_text()).strip()
                        lines_list = page_text.split('\n')
                        lines_list = ['\n'+line.strip()+'\n' if self.__isTitle(line) else line.strip()+' ' for line in lines_list if len(line.strip()) > 0]
                        page_text = ''.join(lines_list).strip()
                        page_text = self.__formatPageText(page_text=page_text)
                        page_number = page_index+1
                        if page_text.startswith(str(page_index)): page_text = page_text[len(str(page_index)):]
                        if main_page != page_number: page_text = self.getSummaryText(text=page_text.strip(), max_tokens=max_tokens)
                        page_text = self.__formatPageText(page_text=page_text)
                        page_text = page_text.strip()
                        if len(page_text) > 0: document_text += f'PAGE {page_number}\n{page_text}\n\n'
                        else: document_text += f'PAGE {page_number}\nBlank page.\n\n'
                    except: document_text += f'PAGE {page_index+1}\nInaccessible content.\n\n'
            else:
                for page_index in range(pages_number):
                    try:
                        page_reader = pdf_reader.pages[page_index]
                        page_text = str(page_reader.extract_text()).strip()
                        lines_list = page_text.split('\n')
                        lines_list = ['\n'+line.strip()+'\n' if self.__isTitle(line) else line.strip()+' ' for line in lines_list if len(line.strip()) > 0]
                        page_text = ''.join(lines_list).strip()
                        page_text = self.__formatPageText(page_text=page_text)
                        if page_text.startswith(str(page_index)): page_text = page_text[len(str(page_index)):]
                        page_text = page_text.strip()
                        if len(page_text) > 0: document_text += page_text+'\n\n'
                    except: pass
                document_text = self.getSummaryText(text=document_text.strip(), max_tokens=original_max_tokens)
                document_text = self.__formatPageText(page_text=document_text)
            if 'pdf_file' in locals(): pdf_file.close()
            tokens_number = self.__countTokens(string=document_text)
            if tokens_number > original_max_tokens: document_text = self.__getBeginningMiddleAndEnd(string=document_text, max_tokens=original_max_tokens)
            result_text = file_type_description+document_text.strip()
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryPDF: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryWord(self, file_path='', max_tokens=1, main_page=None):
        try:
            result_text = ''
            file_path, document_text = str(file_path).strip(), ''
            original_max_tokens = max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            main_page = max(1, int(main_page)) if type(main_page) in (int, float) else None
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                docx_reader = self.__Document(self.__BytesIO(self.__get(file_path).content))
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                docx_reader = self.__Document(file_path)
            paragraphs_number, standard_number_of_paragraphs = len(docx_reader.paragraphs), 16
            pages_number = max((1, int(paragraphs_number/standard_number_of_paragraphs)+1))
            file_type_description = '*Microsoft Word File*\n\n'
            if main_page != None:
                tokens_number = 2184
                if tokens_number < max_tokens: max_tokens = max_tokens-tokens_number
            max_tokens, page_number = int(((max_tokens-len(file_type_description))/pages_number)-len(f'PAGE {pages_number}')), 1
            if original_max_tokens > paragraphs_number and max_tokens > 0:
                for paragraph_index in range(paragraphs_number):
                    try:
                        page_reader = docx_reader.paragraphs[paragraph_index]
                        page_text = str(page_reader.text).strip()
                        lines_list = page_text.split('\n')
                        lines_list = ['\n'+line.strip()+'\n' if self.__isTitle(line) else line.strip()+' ' for line in lines_list if len(line.strip()) > 0]
                        page_text = ''.join(lines_list).strip()
                        page_text = self.__formatPageText(page_text=page_text)
                        paragraph_number = paragraph_index+1
                        if page_text.startswith(str(paragraph_index)): page_text = page_text[len(str(paragraph_index)):]
                        page_text = self.__formatPageText(page_text=page_text)
                        page_text = page_text.strip()
                        text_page_length = len(page_text)
                        if paragraph_number == 1 or paragraph_number % standard_number_of_paragraphs == 0:
                            if text_page_length > 0: document_text += f'\nPAGE {page_number}\n{page_text}\n\n'
                            else: document_text += f'\nPAGE {page_number}\n'
                            page_number += 1
                        elif text_page_length > 0: document_text += page_text+'\n'
                    except: document_text += f'Inaccessible content.\n\n'
                page_list = document_text.split('\nPAGE ')
                page_list = [self.getSummaryText(text=page, max_tokens=max_tokens) if index != main_page else page for index, page in enumerate(page_list)]
                document_text = '\nPAGE '.join(page_list).replace('\n\n', '\n').replace('\nPAGE ', '\n\nPAGE ')
            else:
                for page_index, paragraph in enumerate(docx_reader.paragraphs):
                    try:
                        paragraph_text = str(paragraph.text).strip()
                        lines_list = paragraph_text.split('\n')
                        lines_list = ['\n'+line.strip()+'\n' if self.__isTitle(line) else line.strip()+' ' for line in lines_list if len(line.strip()) > 0]
                        paragraph_text = ''.join(lines_list).strip()
                        paragraph_text = self.__formatPageText(page_text=paragraph_text)
                        if paragraph_text.startswith(str(page_index)): paragraph_text = paragraph_text[len(str(page_index)):]
                        paragraph_text = paragraph_text.strip()
                        if len(paragraph_text) > 0: document_text += paragraph_text+'\n\n'
                    except: pass
                document_text = self.getSummaryText(text=document_text.strip(), max_tokens=original_max_tokens)
                document_text = self.__formatPageText(page_text=document_text)
            tokens_number = self.__countTokens(string=document_text)
            if tokens_number > original_max_tokens: document_text = self.__getBeginningMiddleAndEnd(string=document_text, max_tokens=original_max_tokens)
            result_text = file_type_description+document_text.strip()
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryWord: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryPowerPoint(self, file_path='', max_tokens=1, main_page=None):
        try:
            result_text = ''
            file_path, document_text = str(file_path).strip(), ''
            original_max_tokens = max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            main_page = max(1, int(main_page)) if type(main_page) in (int, float) else None
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                pptx_reader = self.__Presentation(self.__BytesIO(self.__get(file_path).content))
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                pptx_reader = self.__Presentation(file_path)
            pages_number = len(pptx_reader.slides)
            file_type_description = '*Microsoft PowerPoint File*\n\n'
            if main_page != None:
                if main_page > pages_number: main_page = pages_number
                main_page_text = ''
                page_reader = pptx_reader.slides[main_page-1]
                for shape in page_reader.shapes:
                    if hasattr(shape, 'text'): main_page_text += str(shape.text).strip()
                tokens_number = self.__countTokens(string=main_page_text)
                if tokens_number < max_tokens: max_tokens = max_tokens-tokens_number
            max_tokens = int(((max_tokens-len(file_type_description))/pages_number)-len(f'SLIDE {pages_number}'))
            if original_max_tokens > pages_number and max_tokens > 0:
                for page_index in range(pages_number):
                    try:
                        page_reader, page_text = pptx_reader.slides[page_index], ''
                        for shape in page_reader.shapes:
                            if hasattr(shape, 'text'): page_text += str(shape.text).strip()
                        lines_list = page_text.split('\n')
                        lines_list = ['\n'+line.strip()+'\n' if self.__isTitle(line) else line.strip()+' ' for line in lines_list if len(line.strip()) > 0]
                        page_text = ''.join(lines_list).strip()
                        page_text = self.__formatPageText(page_text=page_text)
                        page_number = page_index+1
                        if page_text.startswith(str(page_index)): page_text = page_text[len(str(page_index)):]
                        if main_page != page_number: page_text = self.getSummaryText(text=page_text.strip(), max_tokens=max_tokens)
                        page_text = self.__formatPageText(page_text=page_text)
                        page_text = page_text.strip()
                        if len(page_text) > 0: document_text += f'SLIDE {page_number}\n{page_text}\n\n'
                        else: document_text += f'SLIDE {page_number}\nBlank slide.\n\n'
                    except: document_text += f'SLIDE {page_index+1}\nInaccessible content.\n\n'
            else:
                for page_index in range(pages_number):
                    try:
                        page_reader, page_text = pptx_reader.slides[page_index], ''
                        for shape in page_reader.shapes:
                            if hasattr(shape, 'text'): page_text += str(shape.text).strip()
                        lines_list = page_text.split('\n')
                        lines_list = ['\n'+line.strip()+'\n' if self.__isTitle(line) else line.strip()+' ' for line in lines_list if len(line.strip()) > 0]
                        page_text = ''.join(lines_list).strip()
                        page_text = self.__formatPageText(page_text=page_text)
                        if page_text.startswith(str(page_index)): page_text = page_text[len(str(page_index)):]
                        page_text = page_text.strip()
                        if len(page_text) > 0: document_text += page_text+'\n\n'
                    except: pass
                document_text = self.getSummaryText(text=document_text.strip(), max_tokens=original_max_tokens)
                document_text = self.__formatPageText(page_text=document_text)
            tokens_number = self.__countTokens(string=document_text)
            if tokens_number > original_max_tokens: document_text = self.__getBeginningMiddleAndEnd(string=document_text, max_tokens=original_max_tokens)
            result_text = file_type_description+document_text.strip()
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryPowerPoint: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryCSV(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path = str(file_path).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                csv_reader = self.__read_csv(file_path)
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                csv_reader = self.__read_csv(file_path)
            lines_number = len(csv_reader)
            result_text = f'# CSV File\nTotal number of registration lines: {lines_number}\n'
            column_names, column_values, column_metrics, metrics = [], [], [], '\n## COLUMNS METRICS\n'
            for column_name in csv_reader:
                column_names.append(column_name)
                series = csv_reader[column_name].tolist()
                column_values.append(series)
                metrics += f'### {column_name}\n'
                metrics += f'**Column type:** {self.__getColumnType(vector=series)}\n'
                least_repeated_value = self.__getRepetitionOfValues(vector=series, function=min)
                most_repeated_value = self.__getRepetitionOfValues(vector=series, function=max)
                minimum, maximum = self.__getExtremityOfValues(vector=series)
                average = self.__getAverageOfValues(vector=series)
                median = self.__getMedianOfValues(vector=series)
                standard_deviation = self.__getStandardDeviationOfValues(vector=series)
                variance = self.__getVarianceOfValues(vector=series)
                column_metrics.append((least_repeated_value, most_repeated_value, minimum, maximum, median))
                metrics += f'**Least repeated value:** {least_repeated_value}\n'
                metrics += f'**Most repeated value:** {most_repeated_value}\n'
                metrics += f'**Minimum value:** {minimum}\n'
                metrics += f'**Maximum value:** {maximum}\n'
                metrics += f'**Average of values:** {average}\n'
                metrics += f'**Median of values:** {median}\n'
                metrics += f'**Standard deviation of values:** {standard_deviation}\n'
                metrics += f'**Variance of values:** {variance}\n'
            metrics = metrics.rstrip()
            table = list(zip(*column_values))
            table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
            data_table_metrics = table_data+metrics+'\n\n'
            tokens_number = self.__countTokens(string=result_text+data_table_metrics)
            if tokens_number > max_tokens:
                description = 'Main records contained in the spreadsheet:\n'
                table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics))
                table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                data_table_metrics = description+table_data+metrics+'\n\n'
                tokens_number = self.__countTokens(string=result_text+data_table_metrics)
                if tokens_number > max_tokens:
                    table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics), _metrics_indexes=(2, 3, 4))
                    table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                    data_table_metrics = description+table_data+metrics+'\n\n'
                    tokens_number = self.__countTokens(string=result_text+data_table_metrics)
                    if tokens_number > max_tokens:
                        table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics), _metrics_indexes=(2, 3))
                        table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                        data_table_metrics = description+table_data+metrics+'\n\n'
                        tokens_number = self.__countTokens(string=result_text+data_table_metrics)
                        if tokens_number > max_tokens:
                            table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics), _metrics_indexes=[4])
                            table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                            data_table_metrics = description+table_data+metrics+'\n\n'
                            tokens_number = self.__countTokens(string=result_text+data_table_metrics)
                            if tokens_number > max_tokens:
                                description = f'Column titles: {", ".join(column_names)}\n'
                                data_table_metrics = description+metrics+'\n\n'
            result_text += data_table_metrics
            tokens_number = self.__countTokens(string=result_text)
            if tokens_number > max_tokens: result_text = self.__getBeginningMiddleAndEnd(string=result_text, max_tokens=max_tokens)
            return result_text.strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryCSV: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryExcel(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path = str(file_path).strip()
            original_max_tokens = max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                xlsx_reader = self.__read_excel(self.__BytesIO(self.__get(file_path).content), sheet_name=None)
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                xlsx_reader = self.__read_excel(file_path, sheet_name=None)
            file_type_description = '*Microsoft Excel File*\n\n'
            max_tokens = max((1, int(round((max_tokens-len(file_type_description))/max((1, len(xlsx_reader.items())))))))
            for tab_name, data in xlsx_reader.items():
                blank_rows = data.isnull().all(axis=1)
                blank_columns = data.isnull().all(axis=0)
                if blank_rows.any() or blank_columns.any():
                    data_without_blank_rows = data.loc[~blank_rows]
                    data_without_rows_and_blank_columns = data_without_blank_rows.loc[:,~blank_columns]
                    data_without_rows_and_blank_columns.columns = data_without_rows_and_blank_columns.iloc[0]
                    data_without_rows_and_blank_columns = data_without_rows_and_blank_columns.iloc[1:]
                    xlsx_reader[tab_name] = data_without_rows_and_blank_columns
            for tab_name, data in xlsx_reader.items():
                lines_number = data.shape[0]
                worksheet_header = f'# SPREADSHEET: {tab_name}\nTotal number of registration lines: {lines_number}\n'
                result_text += worksheet_header
                column_names, column_values, column_metrics, metrics = [], [], [], '\n## COLUMNS METRICS\n'
                for column_name, series in data.items():
                    column_names.append(column_name)
                    series = series.tolist()
                    column_values.append(series)
                    metrics += f'### {column_name}\n'
                    metrics += f'**Column type:** {self.__getColumnType(vector=series)}\n'
                    least_repeated_value = self.__getRepetitionOfValues(vector=series, function=min)
                    most_repeated_value = self.__getRepetitionOfValues(vector=series, function=max)
                    minimum, maximum = self.__getExtremityOfValues(vector=series)
                    average = self.__getAverageOfValues(vector=series)
                    median = self.__getMedianOfValues(vector=series)
                    standard_deviation = self.__getStandardDeviationOfValues(vector=series)
                    variance = self.__getVarianceOfValues(vector=series)
                    column_metrics.append((least_repeated_value, most_repeated_value, minimum, maximum, median))
                    metrics += f'**Least repeated value:** {least_repeated_value}\n'
                    metrics += f'**Most repeated value:** {most_repeated_value}\n'
                    metrics += f'**Minimum value:** {minimum}\n'
                    metrics += f'**Maximum value:** {maximum}\n'
                    metrics += f'**Average of values:** {average}\n'
                    metrics += f'**Median of values:** {median}\n'
                    metrics += f'**Standard deviation of values:** {standard_deviation}\n'
                    metrics += f'**Variance of values:** {variance}\n'
                metrics = metrics.rstrip()
                table = list(zip(*column_values))
                table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                data_table_metrics = table_data+metrics+'\n\n'
                tokens_number = self.__countTokens(string=worksheet_header+data_table_metrics)
                if tokens_number > max_tokens:
                    description = 'Main records contained in the spreadsheet:\n'
                    table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics))
                    table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                    data_table_metrics = description+table_data+metrics+'\n\n'
                    tokens_number = self.__countTokens(string=worksheet_header+data_table_metrics)
                    if tokens_number > max_tokens:
                        table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics), _metrics_indexes=(2, 3, 4))
                        table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                        data_table_metrics = description+table_data+metrics+'\n\n'
                        tokens_number = self.__countTokens(string=worksheet_header+data_table_metrics)
                        if tokens_number > max_tokens:
                            table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics), _metrics_indexes=(2, 3))
                            table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                            data_table_metrics = description+table_data+metrics+'\n\n'
                            tokens_number = self.__countTokens(string=worksheet_header+data_table_metrics)
                            if tokens_number > max_tokens:
                                table = self.__getTableSummary(_table=table, _column_metrics=self.__deepcopy(column_metrics), _metrics_indexes=[4])
                                table_data = self.__tabulate(table, headers=column_names, tablefmt='github')
                                data_table_metrics = description+table_data+metrics+'\n\n'
                                tokens_number = self.__countTokens(string=worksheet_header+data_table_metrics)
                                if tokens_number > max_tokens:
                                    description = f'Column titles: {", ".join(column_names)}\n'
                                    data_table_metrics = description+metrics+'\n\n'
                result_text += data_table_metrics
            result_text = file_type_description+result_text
            tokens_number = self.__countTokens(string=result_text)
            if tokens_number > original_max_tokens: result_text = self.__getBeginningMiddleAndEnd(string=result_text, max_tokens=original_max_tokens)
            return result_text.strip()
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryExcel: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryImage(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path = str(file_path).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            if not file_path.startswith('https://') and not file_path.startswith('http://') and not self.__path.exists(file_path):
                print(f'The path to the "{file_path}" file does not exist.')
                return ''
            try:
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('easyocr').setLevel(self.__ERROR)
                except: pass
                use_gpu, textual_elements = self.__cuda.is_available(), []
                reader, image_type = self.__Reader(self.__languages[:6], gpu=use_gpu, verbose=False), 'WEB'
                if file_path.startswith('https://') or file_path.startswith('http://'):
                    temporary = self.__NamedTemporaryFile(delete=False)
                    response = self.__get(file_path)
                    temporary.write(response.content)
                    temporary.close()
                    textual_elements = reader.readtext(temporary.name, detail=0)
                    self.__unlink(temporary.name)
                else: textual_elements, image_type = reader.readtext(file_path, detail=0), self.getFileType(file_path=file_path).upper()
                complete_text = ' '.join(textual_elements) if len(textual_elements) > 0 else ''
                complete_text = complete_text.strip()
            except Exception as error:
                print(error)
                image_type, complete_text = self.getFileType(file_path=file_path).upper(), ''
            def getListOfObjects(file_path=''):
                try:
                    self.__environ['TORCH_HOME'] = './'
                    self.__getLogger("transformers.file_utils").setLevel(self.__ERROR)
                    self.__filterwarnings("ignore")
                    model = self.__models.detection.fasterrcnn_resnet50_fpn(pretrained=True)
                    model.eval()
                    transform = self.__transforms.Compose([self.__transforms.ToTensor()])
                    if file_path.startswith('https://') or file_path.startswith('http://'):
                        try:
                            self.__environ['SSL_CERT_FILE'] = self.__where()
                            self.__getLogger('requests').setLevel(self.__ERROR)
                        except: pass
                        image = self.__Image.open(self.__BytesIO(self.__get(file_path, stream=False).content))
                    else: image = self.__Image.open(file_path)
                    image = transform(image).unsqueeze(0)
                    with self.__no_grad(): outputs = model(image)
                    labels = outputs[0]['labels']
                    labels_map = self.__object_names
                    list_of_objects = list(set([labels_map[label.item()].strip() for label in labels]))
                    return ', '.join(list_of_objects)
                except: return ''
            def extractColors(file_path='', n_colors=3):
                try:
                    file_path, n_colors = str(file_path).strip(), 3 if type(n_colors) not in (bool, int, float) else max((1, int(n_colors)))
                    if file_path.startswith('https://') or file_path.startswith('http://'):
                        try:
                            self.__environ['SSL_CERT_FILE'] = self.__where()
                            self.__getLogger('requests').setLevel(self.__ERROR)
                        except: pass
                        response = self.__get(file_path)
                        image = self.__Image.open(self.__BytesIO(response.content))
                    else: image = self.__Image.open(file_path)
                    image = image.resize((image.width // 10, image.height // 10))
                    image_array = self.__array(image)
                    image_array = image_array.reshape((image_array.shape[0] * image_array.shape[1], 3))
                    with self.__catch_warnings():
                        self.__simplefilter('ignore', category=FutureWarning)
                        kmeans = self.__KMeans(n_clusters=n_colors, n_init='auto')
                        kmeans.fit(image_array)
                    colors = kmeans.cluster_centers_.astype(int)
                    def getColorName(rgb_color=()):
                        def closestColor(requested_color=()):
                            min_colors = {}
                            for key, name in self.__CSS3_HEX_TO_NAMES.items():
                                r_c, g_c, b_c = self.__hex_to_rgb(key)
                                rd = (r_c - requested_color[0]) ** 2
                                gd = (g_c - requested_color[1]) ** 2
                                bd = (b_c - requested_color[2]) ** 2
                                min_colors[(rd + gd + bd)] = name
                            return min_colors[min(min_colors.keys())]
                        try: return self.__rgb_to_name(rgb_color)
                        except ValueError: return closestColor(rgb_color)
                    color_names = list(set([getColorName(tuple(color)) for color in colors]))
                    return ', '.join(color_names)
                except: return ''
            object_labels = getListOfObjects(file_path=file_path).strip()
            color_names = extractColors(file_path=file_path, n_colors=5).strip()
            len_complete_text, len_object_labels, len_color_names = len(complete_text), len(object_labels), len(color_names)
            if sum((len_complete_text, len_object_labels, len_color_names)) > 0:
                list_of_descriptions = []
                if len_complete_text > 0: list_of_descriptions.append('texts = '+complete_text)
                if len_object_labels > 1: list_of_descriptions.append('objects = '+object_labels)
                if len_color_names > 1: list_of_descriptions.append('colors = '+color_names)
                result_text = ' - '.join(list_of_descriptions)
                result_text = image_type+' Image content: '+result_text
            else: result_text = 'The content of the image was not recognized.'
            tokens_number = self.__countTokens(string=result_text)
            if tokens_number > max_tokens: result_text = self.__getBeginningMiddleAndEnd(string=result_text, max_tokens=max_tokens)
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryImage: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryAudio(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path = str(file_path).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            recognizer, audio_type = self.__Recognizer(), 'WAV'
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                web_data = self.__BytesIO(self.__get(file_path).content)
                web_data.seek(0)
                with self.__AudioFile(web_data) as source: audio = recognizer.record(source)
                audio_type = 'WEB'
            else:
                if not self.__path.exists(file_path):
                    print(f'The path to the "{file_path}" file does not exist.')
                    return ''
                if file_path.lower().endswith('.mp3'):
                    audio = self.__AudioSegment.from_mp3(file_path)
                    wav_data = self.__BytesIO()
                    audio.export(wav_data, format='wav')
                    wav_data.seek(0)
                    with self.__AudioFile(wav_data) as source: audio = recognizer.record(source)
                    audio_type = 'MP3'
                else:
                    with self.__AudioFile(file_path) as source: audio = recognizer.record(source)
            def getTranscript(language=''): return recognizer.recognize_google(audio, language=str(language).strip())
            try: temporary_result_text = str(getTranscript(language=self.__languages[0])).strip()
            except: temporary_result_text = ''
            for language in self.__languages[1:]:
                try:
                    result_text = str(getTranscript(language=language)).strip()
                    if len(result_text) > len(temporary_result_text): break
                except: pass
            result_text = result_text.strip()
            result_text = f'{audio_type} audio file content: {result_text}'
            tokens_number = self.__countTokens(string=result_text)
            if tokens_number > max_tokens: result_text = self.__getBeginningMiddleAndEnd(string=result_text, max_tokens=max_tokens)
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryAudio: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryVideo(self, file_path='', max_tokens=1):
        try:
            result_text = ''
            file_path = str(file_path).strip()
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            recognizer, temporary, used_temporary_file = self.__Recognizer(), None, False
            if file_path.startswith('https://') or file_path.startswith('http://'):
                try:
                    self.__environ['SSL_CERT_FILE'] = self.__where()
                    self.__getLogger('requests').setLevel(self.__ERROR)
                except: pass
                response = self.__get(file_path, stream=True)
                if response.status_code == 200:
                    temporary = self.__NamedTemporaryFile(delete=False)
                    temporary.write(response.content)
                    temporary.close()
                    file_path = temporary.name
                    used_temporary_file = True
            elif not self.__path.exists(file_path):
                print(f'The path to the "{file_path}" file does not exist.')
                return ''
            audio_segment = self.__AudioSegment.from_file(file_path)
            wav_bytes = self.__BytesIO()
            audio_segment.export(wav_bytes, format='wav')
            wav_bytes.seek(0)
            if used_temporary_file: self.__unlink(temporary.name)
            with self.__AudioFile(wav_bytes) as source: audio = recognizer.record(source)
            def getTranscript(language=''): return recognizer.recognize_google(audio, language=str(language).strip())
            temporary_result_text = str(getTranscript(language=self.__languages[0])).strip()
            for language in self.__languages[1:]:
                try:
                    result_text = str(getTranscript(language=language)).strip()
                    if len(result_text) > len(temporary_result_text): break
                except: pass
            result_text = result_text.strip()
            result_text = 'In the video the following is said: '+result_text
            tokens_number = self.__countTokens(string=result_text)
            if tokens_number > max_tokens: result_text = self.__getBeginningMiddleAndEnd(string=result_text, max_tokens=max_tokens)
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryVideo: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def getSummaryFile(self, file_path='', max_tokens=1, main_page=None):
        try:
            result_text = ''
            file_path, file_text = str(file_path).strip(), ''
            max_tokens = max((1, int(max_tokens))) if type(max_tokens) in (int, float) else 1
            main_page = max(1, int(main_page)) if type(main_page) in (int, float) else None
            file_type = self.getFileType(file_path=file_path).lower()
            if len(file_type) < 1: file_type = 'null'
            file_path_lower = file_path.lower()
            pdf_types, word_types, text_types = 'pdf', 'docx', 'txt'
            powerpoint_types = ('pptx', 'ppsx', 'pptm')
            excel_types, csv_types = 'xlsx', 'csv'
            image_types = ('webp', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'dng', 'mpo', 'tif', 'tiff', 'pfm')
            audio_types = ('mp3', 'wav', 'mpeg', 'm4a', 'aac', 'ogg', 'flac', 'aiff', 'wma', 'ac3', 'amr')
            video_types = ('mp4', 'avi', 'mkv', 'mov', 'webm', 'flv', '3gp', 'wmv', 'ogv')
            if file_type == pdf_types: file_text = self.getSummaryPDF(file_path=file_path, max_tokens=max_tokens, main_page=main_page)
            elif file_type == word_types: file_text = self.getSummaryWord(file_path=file_path, max_tokens=max_tokens, main_page=main_page)
            elif file_type in ' '.join(powerpoint_types): file_text = self.getSummaryPowerPoint(file_path=file_path, max_tokens=max_tokens, main_page=main_page)
            elif file_type == excel_types: file_text = self.getSummaryExcel(file_path=file_path, max_tokens=max_tokens)
            elif file_type == csv_types: file_text = self.getSummaryCSV(file_path=file_path, max_tokens=max_tokens)
            elif file_type in ' '.join(image_types): file_text = self.getSummaryImage(file_path=file_path, max_tokens=max_tokens)
            elif file_type in ' '.join(audio_types): file_text = self.getSummaryAudio(file_path=file_path, max_tokens=max_tokens)
            elif file_type in ' '.join(video_types): file_text = self.getSummaryVideo(file_path=file_path, max_tokens=max_tokens)
            elif file_path_lower.startswith('https://') or file_path_lower.startswith('http://') or file_path_lower.startswith('www.'):
                has_pdf_type = self.__search(r'\b(' + pdf_types + r')\b', file_path_lower)
                has_word_type = self.__search(r'\b(' + word_types + r')\b', file_path_lower)
                has_powerpoint_type = self.__search(r'\b(' + '|'.join(powerpoint_types) + r')\b', file_path_lower)
                has_excel_type = self.__search(r'\b(' + excel_types + r')\b', file_path_lower)
                has_csv_type = self.__search(r'\b(' + csv_types + r')\b', file_path_lower)
                has_image_type = self.__search(r'\b(' + '|'.join(image_types) + r')\b', file_path_lower)
                has_audio_type = self.__search(r'\b(' + '|'.join(audio_types) + r')\b', file_path_lower)
                has_video_type = self.__search(r'\b(' + '|'.join(video_types) + r')\b', file_path_lower)
                has_text_type = self.__search(r'\b(' + text_types + r')\b', file_path_lower)
                if file_type in 'html php': file_text = self.getSummaryWEBPage(file_path=file_path, max_tokens=max_tokens)
                elif 'youtube.com' in file_path_lower or 'youtu.be' in file_path_lower: file_text = self.getSummaryYouTube(file_path=file_path, max_tokens=max_tokens)
                elif 'image' in file_path_lower or '/img/' in file_path_lower or has_image_type: file_text = self.getSummaryImage(file_path=file_path, max_tokens=max_tokens)
                elif has_pdf_type: file_text = self.getSummaryPDF(file_path=file_path, max_tokens=max_tokens, main_page=main_page)
                elif has_word_type: file_text = self.getSummaryWord(file_path=file_path, max_tokens=max_tokens, main_page=main_page)
                elif has_powerpoint_type: file_text = self.getSummaryPowerPoint(file_path=file_path, max_tokens=max_tokens, main_page=main_page)
                elif has_excel_type: file_text = self.getSummaryExcel(file_path=file_path, max_tokens=max_tokens)
                elif has_csv_type: file_text = self.getSummaryCSV(file_path=file_path, max_tokens=max_tokens)
                elif has_audio_type: file_text = self.getSummaryAudio(file_path=file_path, max_tokens=max_tokens)
                elif has_video_type: file_text = self.getSummaryVideo(file_path=file_path, max_tokens=max_tokens)
                elif has_text_type: file_text = self.getSummaryTXT(file_path=file_path, max_tokens=max_tokens)
                else: file_text = self.getSummaryWEBPage(file_path=file_path, max_tokens=max_tokens)
            else: file_text = self.getSummaryTXT(file_path=file_path, max_tokens=max_tokens)
            tokens_number = self.__countTokens(string=file_text)
            if tokens_number > max_tokens: file_text = self.__getBeginningAndEnd(string=file_text, max_tokens=max_tokens)
            result_text = file_text.strip()
            return result_text
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getSummaryFile: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return ''
    def saveContext(self, user_id=0, dialog_id=0, prompt='', answer=''):
        try:
            total_tokens = self.__countTokens(string=prompt)+self.__countTokens(string=answer)
            user_id, dialog_id, prompt, answer = str(user_id).strip(), str(dialog_id).strip(), str(prompt).strip(), str(answer).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            if not self.__path.isdir(context_directory): self.__makedirs(context_directory)
            user_context = f'{context_directory}{user_id}/'
            if not self.__path.isdir(user_context): self.__makedirs(user_context)
            all_tokens = total_tokens+self.__getAllTokens(user_id=user_id, dialog_id=dialog_id)
            self.__saveAllTokens(user_id=user_id, dialog_id=dialog_id, all_tokens=all_tokens)
            dialogue_context = f'{user_context}{dialog_id}/'
            if not self.__path.isdir(dialogue_context): self.__makedirs(dialogue_context)
            conversation_id = len(self.__listdir(dialogue_context))
            conversation_context = f'{dialogue_context}{conversation_id}.index'
            if not self.__path.exists(conversation_context):
                write = open(conversation_context, 'w', encoding='utf-8', errors='ignore')
                conversation = f'{prompt}\n<|conversation|>\n{answer}'
                write.write(self.__utf8ToBase64(utf8_string=conversation))
                write.close()
            return True
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.saveContext: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return False
    def deleteContext(self, user_id=0, dialog_id=0):
        try:
            user_id, dialog_id = str(user_id).strip(), str(dialog_id).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            user_context = f'{context_directory}{user_id}/'
            dialogue_context = f'{user_context}{dialog_id}/'
            tokens_context = f'{user_context}{dialog_id}.tokens'
            if self.__path.isdir(dialogue_context): self.__rmtree(dialogue_context)
            if self.__path.exists(tokens_context): self.__unlink(tokens_context)
            return True if not self.__path.isdir(dialogue_context) and not self.__path.exists(tokens_context) else False
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.deleteContext: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return False
    def getContext(self, user_id=0, dialog_id=0, config=None):
        try:
            user_id, dialog_id = str(user_id).strip(), str(dialog_id).strip()
            user_id, dialog_id = user_id if len(user_id) > 0 else '0', dialog_id if len(dialog_id) > 0 else '0'
            default_config = {
                'system': '',
                'prompt': '',
                'max_tokens': 1,
                'return_format': 'dictionaries_list',
                'system_key': 'system',
                'interlocutor_key': 'role',
                'user_value': 'user',
                'assistant_value': 'assistant',
                'content_key': 'content',
                'dialogue_indexes': []
            }
            config = config if type(config) == dict else default_config
            system = str(config['system']).strip() if 'system' in config else ''
            prompt = str(config['prompt']).strip() if 'prompt' in config else ''
            system_tokens, prompt_tokens = self.__countTokens(string=system), self.__countTokens(string=prompt)
            max_tokens = max((1, int(config['max_tokens']))) if 'max_tokens' in config else 1
            google_system = '(SYSTEM) You must respond as if you were the following: '
            return_format = str(config['return_format']).lower().strip() if 'return_format' in config else 'dictionaries_list'
            if return_format in ('gemini_pattern', 'gemma_pattern'): max_tokens = max((1, max_tokens-(system_tokens+prompt_tokens+self.__countTokens(string=google_system))))
            else: max_tokens = max((1, max_tokens-(system_tokens+prompt_tokens)))
            number_of_dialogues = self.__getNumberOfDialogues(user_id=user_id, dialog_id=dialog_id)
            if return_format == 'claude_pattern': max_tokens = max((1, max_tokens-(4+(7*number_of_dialogues)+6)))
            elif return_format == 'llama3_pattern': max_tokens = max((1, max_tokens-(27+(41*number_of_dialogues)+35)))
            elif return_format == 'mistral_pattern': max_tokens = max((1, max_tokens-(6+(10*number_of_dialogues)+6)))
            elif return_format == 'gemma_pattern': max_tokens = max((1, max_tokens-(13+(24*number_of_dialogues)+18)))
            elif return_format == 'phi3_pattern': max_tokens = max((1, max_tokens-(9+(18*number_of_dialogues)+14)))
            elif return_format == 'yi_pattern': max_tokens = max((1, max_tokens-(13+(26*number_of_dialogues)+20)))
            elif return_format == 'falcon_pattern': max_tokens = max((1, max_tokens-(1+(6*number_of_dialogues)+5)))
            elif return_format == 'falcon2_pattern': max_tokens = max((1, max_tokens-(3+(7*number_of_dialogues)+6)))
            elif return_format == 'stablelm2_pattern': max_tokens = max((1, max_tokens-(13+(26*number_of_dialogues)+20)))
            system_key = str(config['system_key']).strip() if 'system_key' in config else 'system'
            interlocutor_key = str(config['interlocutor_key']).strip() if 'interlocutor_key' in config else 'role'
            user_value = str(config['user_value']).strip() if 'user_value' in config else 'user'
            assistant_value = str(config['assistant_value']).strip() if 'assistant_value' in config else 'assistant'
            content_key = str(config['content_key']).strip() if 'content_key' in config else 'content'
            dialogue_indexes = config['dialogue_indexes'] if 'dialogue_indexes' in config else []
            dialogue_indexes = list(dialogue_indexes) if type(dialogue_indexes) in (tuple, list) else []
            all_tokens = self.__getAllTokens(user_id=user_id, dialog_id=dialog_id)
            context_directory = f'{self.__getRootDirectory()}context_directory/'
            user_context = f'{context_directory}{user_id}/'
            dialogue_context = f'{user_context}{dialog_id}/'
            if self.__path.isdir(dialogue_context):
                prompt_keywords, context_window = self.__getKeyWords(string=prompt), []
                if all_tokens <= max_tokens:
                    if len(dialogue_indexes) > 0: context_window = self.__getRelatedContext(user_id=user_id, dialog_id=dialog_id, key_words=prompt_keywords, max_tokens=max_tokens, indexes=dialogue_indexes)
                    else:
                        dialog_length = len(self.__listdir(dialogue_context))
                        for index in range(0, dialog_length): context_window.append(self.__getIndex(user_id=user_id, dialog_id=dialog_id, index=index))
                else: context_window = self.__getRelatedContext(user_id=user_id, dialog_id=dialog_id, key_words=prompt_keywords, max_tokens=max_tokens, indexes=dialogue_indexes)
                context_tokens = self.__countTokens(string=str(context_window))
                number_of_tokens = system_tokens+prompt_tokens+context_tokens
                if number_of_tokens > max_tokens:
                    system, prompt = self.__getSummaryText(text=system, max_tokens=max_tokens), self.__getSummaryText(text=prompt, max_tokens=max_tokens)
                    system_tokens, prompt_tokens = self.__countTokens(string=system), self.__countTokens(string=prompt)
                    max_tokens = max_tokens-system_tokens-prompt_tokens
                    max_tokens = max((1, max_tokens/max((1, len(context_window)))))
                    context_window = self.__getRelatedContext(user_id=user_id, dialog_id=dialog_id, key_words=prompt_keywords, max_tokens=max_tokens, indexes=dialogue_indexes)
                if return_format == 'dictionaries_list':
                    messages = [{interlocutor_key: system_key, content_key: system}] if len(system) > 0 else []
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        messages.append({interlocutor_key: user_value, content_key: _input.strip()}), messages.append({interlocutor_key: assistant_value, content_key: _output.strip()})
                    if len(prompt) > 0: messages.append({interlocutor_key: user_value, content_key: prompt})
                    return messages
                elif return_format == 'chatgpt_pattern':
                    messages = [{"role": "system", "content": system}] if len(system) > 0 else []
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        messages.append({"role": "user", "content": _input.strip()}), messages.append({"role": "assistant", "content": _output.strip()})
                    if len(prompt) > 0: messages.append({"role": "user", "content": prompt})
                    return messages
                elif return_format == 'gemini_pattern':
                    contents = [{"role": "user", "parts": [{"text": google_system+system}]}] if len(system) > 0 else []
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        contents.append({"role": "user", "parts": [{"text": _input.strip()}]}), contents.append({"role": "model", "parts": [{"text": _output.strip()}]})
                    if len(prompt) > 0: contents.append({"role": "user", "parts": [{"text": prompt}]})
                    return contents
                elif return_format == 'claude_pattern':
                    template = f'\n\nSystem: '+system if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'\n\nHuman: '+_input.strip()
                        template += f'\n\nAssistant: '+_output.strip()
                    if len(prompt) > 0: template += f'\n\nHuman: {prompt}\n\nAssistant:'
                    return template
                elif return_format == 'llama3_pattern':
                    template = f'<|begin_of_text|><|start_header_id|>system<|end_header_id|>\n\n{system}<|eot_id|>' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'<|start_header_id|>user<|end_header_id|>\n\n{_input.strip()}<|eot_id|>'
                        template += f'<|start_header_id|>assistant<|end_header_id|>\n\n{_output.strip()}<|eot_id|>'
                    if len(prompt) > 0: template += f'<|start_header_id|>user<|end_header_id|>\n\n{prompt}<|eot_id|><|start_header_id|>assistant<|end_header_id|>\n\n'
                    return template
                elif return_format in ('mistral_pattern', 'mixtral_pattern'):
                    template = f'[INST]{system}[/INST]' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'<s>[INST]{_input.strip()}[/INST]{_output.strip()}</s>'
                    if len(prompt) > 0: template += f'[INST]{prompt}[/INST]'
                    return template
                elif return_format in ('gemma_pattern', 'gemma2_pattern'):
                    template = f'<start_of_turn>user\n{google_system}{system}<end_of_turn>\n' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'<start_of_turn>user\n{_input.strip()}<end_of_turn>\n'
                        template += f'<start_of_turn>model\n{_input.strip()}<end_of_turn>\n'
                    if len(prompt) > 0: template += f'<start_of_turn>user\n{prompt}<end_of_turn>\n<start_of_turn>model'
                    return template
                elif return_format == 'phi3_pattern':
                    template = f'<|system|>\n{system}<|end|>\n' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'<|user|>\n{_input.strip()}<|end|>\n'
                        template += f'<|assistant|>\n{_output.strip()}<|end|>\n'
                    if len(prompt) > 0: template += f'<|user|>\n{prompt}<|end|>\n<|assistant|>'
                    return template
                elif return_format == 'yi_pattern':
                    template = f'<|im_start|>system\n{system}<|im_end|>\n' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'<|im_start|>user\n{_input.strip()}<|im_end|>\n'
                        template += f'<|im_start|>assistant\n{_output.strip()}<|im_end|>\n'
                    if len(prompt) > 0: template += f'<|im_start|>user\n{prompt}<|im_end|>\n<|im_start|>assistant'
                    return template
                elif return_format == 'falcon_pattern':
                    template = system+'\n' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'User: {_input.strip()}\n'
                        template += f'Assistant: {_output.strip()}\n'
                    if len(prompt) > 0: template += f'User: {prompt}\nAssistant:'
                    return template
                elif return_format == 'falcon2_pattern':
                    template = f'System: {system}\n' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'User: {_input.strip()}\n'
                        template += f'Falcon: {_output.strip()}\n'
                    if len(prompt) > 0: template += f'User: {prompt}\nFalcon:'
                    return template
                elif return_format == 'stablelm2_pattern':
                    template = f'<|im_start|>system\n{system}<|im_end|>\n' if len(system) > 0 else ''
                    for context in context_window:
                        _input, _output = context.split('<|conversation|>')
                        template += f'<|im_start|>user\n{_input.strip()}<|im_end|>\n'
                        template += f'<|im_start|>assistant\n{_output.strip()}<|im_end|>\n'
                    if len(prompt) > 0: template += f'<|im_start|>user\n{prompt}<|im_end|>\n<|im_start|>assistant'
                    return template
                else:
                    infinite_context = '\n<|context|>\n'.join(context_window)
                    prompt_content = '\n<|prompt|>\n'+prompt+'\n<|prompt|>\n' if len(prompt) > 0 else ''
                    system_content = '\n<|system|>\n'+system+'\n<|system|>\n' if len(system) > 0 else ''
                    infinite_context = system_content+infinite_context+prompt_content
                    return infinite_context.strip()
            else: return prompt if prompt_tokens <= max_tokens else self.__getSummaryText(text=prompt, max_tokens=max_tokens)
        except Exception as error:
            error_message = 'ERROR in PerpetualContext.getContext: '+str(error)
            print(error_message)
            try: self.__print_exc() if self.__display_error_point else None
            except: pass
            return prompt
# This code is an algorithm projected, architected and developed by Sapiens Technology®️ and aims to manage context window memory
# in Artificial Intelligence projects for language models. It manages context memory by saving and indexing the encoded dialogs
# for later consultation and return of excerpts referring to the input prompt, summarizing these excerpts when necessary to prevent
# the character sequence from exceeding the tokens limit previously established by the control variable.
# This makes it possible to establish a perpetual and limitless context even for language models with limited context windows.
